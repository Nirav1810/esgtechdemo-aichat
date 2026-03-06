"""
generate_from_template.py
─────────────────────────
Standalone PPTX generator that uses a Slidesgo template as the base
presentation and populates it according to a JSON slide plan.

Usage
─────
# Inspect the template to find layout indices & placeholder slots:
    python generate_from_template.py --inspect-template "Sustainability and Environment Newsletter _ by Slidesgo.pptx"

# Generate a presentation from a slide plan:
    python generate_from_template.py plan.json

# Override template or output path at runtime:
    python generate_from_template.py plan.json --template "My Template.pptx" --output "report.pptx"

Dependencies
────────────
    pip install python-pptx matplotlib seaborn requests pillow

Optional (for Pexels stock photos):
    Set PEXELS_API_KEY in your environment or in a .env file.
"""

from __future__ import annotations

import argparse
import json
import math
import os
import sys
import warnings
from io import BytesIO
from pathlib import Path
from typing import Any

import matplotlib
matplotlib.use("Agg")  # Non-interactive backend — must be set before pyplot import
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
import requests
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.enum.dml import MSO_THEME_COLOR

# ─────────────────────────────────────────────
# Brand / style constants
# ─────────────────────────────────────────────

BRAND_PALETTE = [
    "#1B4332",  # deep green
    "#2D6A4F",  # forest green
    "#40916C",  # medium green
    "#52B788",  # sage green
    "#74C69D",  # light green
    "#95D5B2",  # mint
    "#B7E4C7",  # pale green
]

ACCENT_BLUE = "#1565C0"
TEXT_DARK   = "#1A1A2E"
TEXT_LIGHT  = "#FFFFFF"

# Matplotlib / Seaborn chart style applied globally
CHART_STYLE = {
    "figure.facecolor":   "#EFF8F1",
    "axes.facecolor":     "#EFF8F1",
    "axes.spines.top":    False,
    "axes.spines.right":  False,
    "axes.spines.left":   True,
    "axes.spines.bottom": True,
    "axes.grid":          True,
    "grid.color":         "#C8E8D0",
    "grid.linewidth":     0.8,
    "font.family":        "DejaVu Sans",
    "font.size":          11,
    "axes.labelsize":     12,
    "axes.titlesize":     14,
    "xtick.labelsize":    11,
    "ytick.labelsize":    11,
    "legend.fontsize":    11,
    "figure.dpi":         150,
}

CHART_DPI = 150  # Embedded chart resolution
MIN_TEXT_PT = 12

# ─────────────────────────────────────────────
# Helper: .env loader (lightweight, no dotenv dep)
# ─────────────────────────────────────────────

def _load_dotenv() -> None:
    """
    Load environment variables from .env.local (Next.js convention) or .env.
    Searches the script's directory and each parent up to the filesystem root.
    Does not overwrite variables already present in the environment.
    """
    search_dirs = [Path(__file__).resolve().parent]
    search_dirs += list(search_dirs[0].parents)

    for directory in search_dirs:
        for filename in (".env.local", ".env"):
            env_path = directory / filename
            if env_path.exists():
                for line in env_path.read_text(encoding="utf-8").splitlines():
                    line = line.strip()
                    if not line or line.startswith("#") or "=" not in line:
                        continue
                    key, _, value = line.partition("=")
                    key = key.strip()
                    value = value.strip().strip('"').strip("'")
                    if key and key not in os.environ:
                        os.environ[key] = value
                return  # stop after the first file found

_load_dotenv()

# ─────────────────────────────────────────────
# Template slide removal
# ─────────────────────────────────────────────

# OPC relationship namespace
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _clear_template_slides(prs: Presentation) -> None:
    """
    Remove every sample slide that ships inside the template while keeping
    the slide masters, layouts, theme, and fonts completely intact.

    python-pptx requires TWO steps to fully delete a slide:
      1. Drop the OPC relationship from the presentation part (so the slide
         XML file is no longer referenced and won't appear in the output).
      2. Remove the <p:sldId> element from the sldIdLst so the slide index
         is cleared from the presentation XML.
    Doing only step 2 (the common mistake) leaves orphaned parts that some
    PPTX viewers still render.
    """
    sldIdLst = prs.slides._sldIdLst
    rIds = [
        sld.get(f"{{{_R_NS}}}id")
        for sld in list(sldIdLst)
    ]
    for rId in rIds:
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass  # already gone
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)


# ─────────────────────────────────────────────

def list_layouts(prs: Presentation) -> None:
    """
    Print a table of every slide layout in the template, including:
      - layout index (use this as `layout_index` in your JSON)
      - layout name
      - each placeholder's idx, type, and bounding box
    """
    print(f"\n{'─'*70}")
    print(f"  Template: {prs.core_properties.title or 'Untitled'}")
    print(f"  Slide layouts found: {len(prs.slide_layouts)}")
    print(f"{'─'*70}\n")

    for i, layout in enumerate(prs.slide_layouts):
        print(f"[{i:02d}]  Layout name: \"{layout.name}\"")
        if layout.placeholders:
            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type).split(".")[-1]
                left   = round(ph.left   / 914400, 2) if ph.left   else "?"
                top    = round(ph.top    / 914400, 2) if ph.top    else "?"
                width  = round(ph.width  / 914400, 2) if ph.width  else "?"
                height = round(ph.height / 914400, 2) if ph.height else "?"
                print(f"       ph_idx={ph.placeholder_format.idx:<3}  "
                      f"type={ph_type:<20}  "
                      f"name=\"{ph.name}\"  "
                      f"pos=({left}\", {top}\")  "
                      f"size=({width}\" × {height}\")")
        else:
            print("       (no placeholders)")
        print()

# ─────────────────────────────────────────────
# Layout resolution
# ─────────────────────────────────────────────

def get_layout(prs: Presentation, slide_data: dict[str, Any]):
    """
    Return the slide layout requested in slide_data.
    Priority:
      1. `layout_index` (int)  — direct index into prs.slide_layouts
      2. `layout_name`  (str)  — case-insensitive partial match on layout.name
    Raises ValueError with a helpful list of available layouts on failure.
    """
    layouts = prs.slide_layouts

    # By index
    if "layout_index" in slide_data:
        idx = int(slide_data["layout_index"])
        if 0 <= idx < len(layouts):
            return layouts[idx]
        available = "\n".join(f"  [{i}] {l.name}" for i, l in enumerate(layouts))
        raise ValueError(
            f"layout_index={idx} is out of range (0–{len(layouts)-1}).\n"
            f"Available layouts:\n{available}"
        )

    # By name substring
    if "layout_name" in slide_data:
        needle = slide_data["layout_name"].lower()
        for layout in layouts:
            if needle in layout.name.lower():
                return layout
        available = "\n".join(f"  {l.name}" for l in layouts)
        raise ValueError(
            f"No layout matching \"{slide_data['layout_name']}\".\n"
            f"Available layouts:\n{available}"
        )

    # Fallback: first layout (index 0)
    warnings.warn(
        f"Slide id={slide_data.get('id','?')} has no layout_index or layout_name — "
        "using layout 0 as fallback.",
        stacklevel=2,
    )
    return layouts[0]

# ─────────────────────────────────────────────
# Placeholder helpers
# ─────────────────────────────────────────────

def fill_placeholder(slide, ph_idx: int, text: str, warn_missing: bool = True) -> bool:
    """
    Write `text` into the placeholder with the given idx.
    Clears existing runs first to avoid ghost text from the template.
    Returns True on success, False (with a warning) if not found.
    """
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        if warn_missing:
            warnings.warn(f"WARNING: placeholder idx={ph_idx} not found on this slide layout.")
        return False

    tf = ph.text_frame
    tf.clear()
    tf.text = text
    return True


def set_placeholder_font(slide, ph_idx: int,
                          bold: bool | None = None,
                          size_pt: int | None = None,
                          color_hex: str | None = None) -> None:
    """Apply optional font overrides to the first paragraph of a placeholder."""
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return

    tf = ph.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if bold is not None:
                run.font.bold = bold
            if size_pt is not None:
                run.font.size = Pt(size_pt)
            if color_hex is not None:
                run.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))


def add_bullets_to_placeholder(slide, ph_idx: int, bullets: list[str]) -> bool:
    """
    Fill a text placeholder with a list of bullet strings.
    Uses add_paragraph() for each item after the first.
    Returns False if the placeholder does not exist OR is too small to show
    body text (height < 0.75"), so the caller can fall back to a textbox.
    """
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        warnings.warn(f"WARNING: placeholder idx={ph_idx} not found on this slide layout.")
        return False

    # Guard: skip tiny placeholders (< 0.75") that can't visibly hold body text
    if ph.height and ph.height < Inches(0.75):
        return False

    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
            para.text = bullet
        else:
            para = tf.add_paragraph()
            para.text = bullet
            para.level = 0

        para.line_spacing = 1.3
        for run in para.runs:
            run.font.size = Pt(max(MIN_TEXT_PT, 13))
            run.font.color.rgb = RGBColor.from_string(TEXT_DARK.lstrip("#"))

    return True

# ─────────────────────────────────────────────
# Chart rendering — Matplotlib + Seaborn
# ─────────────────────────────────────────────

def _apply_chart_style() -> None:
    """Apply the global brand chart style to rcParams."""
    plt.rcParams.update(CHART_STYLE)


def _palette(cfg: dict[str, Any], n: int) -> list[str]:
    """Return a list of `n` hex colour strings from the chart config or brand defaults."""
    user = cfg.get("color_palette")
    if user and isinstance(user, list) and len(user) >= n:
        return user[:n]
    # Cycle brand palette
    return [BRAND_PALETTE[i % len(BRAND_PALETTE)] for i in range(n)]


def render_chart_png(chart_cfg: dict[str, Any], fig_aspect: float | None = None) -> BytesIO:
    """
    Render a chart to a BytesIO PNG buffer using Matplotlib / Seaborn.

    Supported chart_type values:
        bar            — vertical bar (single or grouped)
        horizontal_bar — horizontal bar
        stacked_bar    — stacked vertical bar
        line           — line chart (single or multi-series)
        multi_line     — alias for line with multiple datasets
        area           — filled area chart
        pie            — standard pie chart
        doughnut       — pie with a hole in the centre
        scatter        — scatter / bubble plot

    chart_cfg schema:
    {
        "chart_type": "bar",
        "title": "Scope 1 Emissions",          # optional
        "labels": ["Stationary", "Fugitive"],   # x-axis / category labels
        "datasets": [
            {
                "label": "tCO₂Eq",
                "data":  [29746, 5785]
            }
        ],
        "color_palette": ["#2D6A4F", "#52B788"],  # optional
        "x_label": "Source",                      # optional
        "y_label": "tCO₂Eq",                      # optional
        "show_legend": false,                      # optional (default: auto)
        "show_values": true                        # optional: label bars/slices
    }
    """
    _apply_chart_style()

    # (variables already assigned in the sizing block above)

    CHART_BG = CHART_STYLE["figure.facecolor"]  # "#EFF8F1"
    chart_type  = chart_cfg.get("chart_type", "bar")  # read early for sizing
    labels      = chart_cfg.get("labels", [])

    # ── Compute figure size from the display-box aspect ratio ────────────────
    # Caller may pass fig_aspect = display_width / display_height so that the
    # saved PNG has zero stretch distortion when placed in the slide box.
    if fig_aspect is None:
        fig_aspect = 1.6  # safe default (≈ standard 16:10)
    FIG_W = 9.0  # reference width in matplotlib inches (high enough for crisp DPI)
    FIG_H = FIG_W / fig_aspect
    # For horizontal_bar, scale height so each bar gets ~0.9" of figure height
    if chart_type == "horizontal_bar" and labels:
        FIG_H = max(FIG_H, len(labels) * 0.9)
    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H), facecolor=CHART_BG)
    fig.patch.set_facecolor(CHART_BG)
    ax.set_facecolor(CHART_BG)

    # Re-read so the later code still works (chart_type / labels already set above)
    datasets    = chart_cfg.get("datasets", [])
    title       = chart_cfg.get("title", "")
    x_label     = chart_cfg.get("x_label", "")
    y_label     = chart_cfg.get("y_label", "")
    show_values = chart_cfg.get("show_values", True)
    n_series    = max(len(datasets), 1)

    # ── bar ──────────────────────────────────────────────────────────────────
    if chart_type == "bar":
        n_groups = len(labels)
        n_ds     = len(datasets)
        x        = range(n_groups)
        width    = 0.7 / n_ds if n_ds > 1 else 0.55
        palette  = _palette(chart_cfg, n_ds)

        for di, ds in enumerate(datasets):
            offsets = [xi + (di - (n_ds - 1) / 2) * width for xi in x]
            bars = ax.bar(offsets, ds["data"], width=width,
                          color=palette[di], label=ds.get("label", ""),
                          edgecolor="white", linewidth=0.5, zorder=3)
            if show_values:
                for bar in bars:
                    h = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width() / 2, h * 1.01,
                            f"{h:,.0f}", ha="center", va="bottom",
                        fontsize=10, color="#333333")
    # ── horizontal_bar ───────────────────────────────────────────────────────
    elif chart_type == "horizontal_bar":
        palette = _palette(chart_cfg, len(labels))
        data = datasets[0]["data"] if datasets else []
        # Truncate very long y-axis labels so they don't consume the whole figure
        _ylabels = [lb if len(lb) <= 22 else lb[:21] + "\u2026" for lb in labels]
        bars = ax.barh(_ylabels, data, color=palette,
                       edgecolor="white", linewidth=0.5, zorder=3)
        ax.invert_yaxis()
        # Generous left margin so labels don't overlap bars
        fig.subplots_adjust(left=0.30)
        if show_values:
            for bar in bars:
                w = bar.get_width()
                ax.text(w * 1.005, bar.get_y() + bar.get_height() / 2,
                        f"{w:,.0f}", va="center", fontsize=10, color="#333333")

    # ── stacked_bar ──────────────────────────────────────────────────────────
    elif chart_type == "stacked_bar":
        n_ds    = len(datasets)
        palette = _palette(chart_cfg, n_ds)
        x       = range(len(labels))
        bottoms = [0.0] * len(labels)

        for di, ds in enumerate(datasets):
            ax.bar(list(x), ds["data"], bottom=bottoms,
                   color=palette[di], label=ds.get("label", ""),
                   edgecolor="white", linewidth=0.5, zorder=3)
            bottoms = [b + v for b, v in zip(bottoms, ds["data"])]

        ax.set_xticks(list(x))
        _sb_labels = [lb if len(lb) <= 18 else lb[:17] + "\u2026" for lb in labels]
        ax.set_xticklabels(_sb_labels, rotation=30 if len(labels) > 5 else 0, ha="right")

    # ── line / multi_line ─────────────────────────────────────────────────────
    elif chart_type in ("line", "multi_line"):
        palette = _palette(chart_cfg, n_series)
        x_idx = list(range(len(labels)))
        for di, ds in enumerate(datasets):
            ax.plot(x_idx, ds["data"],
                    color=palette[di], label=ds.get("label", ""),
                    linewidth=2.5, marker="o", markersize=5,
                    markeredgecolor="white", markeredgewidth=1)
        ax.set_xticks(x_idx)
        _ll = [lb if len(lb) <= 18 else lb[:17] + "\u2026" for lb in labels]
        ax.set_xticklabels(_ll, rotation=30 if len(labels) > 5 else 0, ha="right")

    # ── area ────────────────────────────────────────────────────────────────
    elif chart_type == "area":
        palette = _palette(chart_cfg, n_series)
        x_idx = list(range(len(labels)))
        for di, ds in enumerate(datasets):
            ax.fill_between(x_idx, ds["data"],
                            alpha=0.35, color=palette[di], label=ds.get("label", ""))
            ax.plot(x_idx, ds["data"],
                    color=palette[di], linewidth=2)
        ax.set_xticks(x_idx)
        _al = [lb if len(lb) <= 18 else lb[:17] + "\u2026" for lb in labels]
        ax.set_xticklabels(_al, rotation=30 if len(labels) > 5 else 0, ha="right")

    # ── pie ──────────────────────────────────────────────────────────────────
    elif chart_type == "pie":
        data    = datasets[0]["data"] if datasets else []
        palette = _palette(chart_cfg, len(data))
        wedges, texts, autotexts = ax.pie(
            data, labels=None, colors=palette,
            autopct="%1.1f%%" if show_values else None,
            startangle=120, pctdistance=0.75,
            wedgeprops=dict(edgecolor="white", linewidth=1.5),
        )
        for at in (autotexts or []):
            at.set_fontsize(10)
            at.set_color("white")
        ax.legend(labels, loc="lower center",
                  fontsize=10, framealpha=0.7,
                  bbox_to_anchor=(0.5, -0.18), ncol=min(len(labels), 3))
        ax.axis("equal")

    # ── doughnut ─────────────────────────────────────────────────────────────
    elif chart_type == "doughnut":
        data    = datasets[0]["data"] if datasets else []
        palette = _palette(chart_cfg, len(data))
        wedges, texts, autotexts = ax.pie(
            data, labels=None, colors=palette,
            autopct="%1.1f%%" if show_values else None,
            startangle=120, pctdistance=0.82,
            wedgeprops=dict(width=0.45, edgecolor="white", linewidth=1.5),
        )
        for at in (autotexts or []):
            at.set_fontsize(10)
            at.set_color("#333333")
        # Centre hole label
        total = sum(data)
        ax.text(0, 0, f"{total:,.0f}", ha="center", va="center",
                fontsize=13, fontweight="bold", color=TEXT_DARK)
        ax.legend(labels, loc="lower center",
                  fontsize=10, framealpha=0.7,
                  bbox_to_anchor=(0.5, -0.18), ncol=min(len(labels), 3))
        ax.axis("equal")

    # ── scatter ───────────────────────────────────────────────────────────────
    elif chart_type == "scatter":
        palette = _palette(chart_cfg, n_series)
        for di, ds in enumerate(datasets):
            x_vals = ds.get("x", list(range(len(ds["data"]))))
            ax.scatter(x_vals, ds["data"], color=palette[di],
                       label=ds.get("label", ""), s=60, alpha=0.8,
                       edgecolors="white", linewidths=0.5)

    else:
        warnings.warn(f"Unknown chart_type='{chart_type}'. Rendering empty chart.")

    # ── shared formatting ────────────────────────────────────────────────────
    if title:
        ax.set_title(title, fontsize=12, fontweight="bold",
                     color=TEXT_DARK, pad=10)
    if x_label and chart_type not in ("pie", "doughnut"):
        ax.set_xlabel(x_label, labelpad=6)
    if y_label and chart_type not in ("pie", "doughnut", "horizontal_bar"):
        ax.set_ylabel(y_label, labelpad=6)
    if chart_type == "horizontal_bar" and x_label:
        ax.set_xlabel(x_label, labelpad=6)

    ax.tick_params(axis="both", colors="#1F4A2E", labelsize=11)
    for spine in ax.spines.values():
        spine.set_color("#9CCFAE")

    # Legend: show only when multiple series or explicitly requested
    show_legend = chart_cfg.get("show_legend")
    if show_legend is True:
        ax.legend(loc="best")
    elif show_legend is None and n_series > 1 and chart_type not in ("pie", "doughnut"):
        ax.legend(loc="best")

    # Pie/doughnut use bottom-centre legend — give extra bottom padding
    if chart_type in ("pie", "doughnut"):
        fig.tight_layout(pad=1.2)
        fig.subplots_adjust(bottom=0.22)
    elif chart_type == "horizontal_bar":
        # left margin already set via subplots_adjust; tight_layout may override it
        fig.tight_layout(pad=1.0)
        fig.subplots_adjust(left=0.30)
    else:
        fig.tight_layout(pad=1.2)

    buf = BytesIO()
    fig.savefig(
        buf,
        format="png",
        dpi=CHART_DPI,
        bbox_inches="tight",
        facecolor=CHART_BG,
        pad_inches=0.15,
    )
    plt.close(fig)
    buf.seek(0)
    return buf

# ─────────────────────────────────────────────
# Image fetching
# ─────────────────────────────────────────────

def fetch_image_bytes(source: str | None) -> BytesIO | None:
    """
    Resolve `source` to a BytesIO PNG/JPEG buffer.
    Priority:
      1. HTTP/HTTPS URL  → download directly
      2. Local file path → read from disk
      3. Keyword string  → Pexels API search (requires PEXELS_API_KEY env var)
    Returns None with a warning if all attempts fail.
    """
    if not source:
        return None

    # ── HTTP URL ────────────────────────────────────────────────────────────
    if source.startswith("http://") or source.startswith("https://"):
        try:
            resp = requests.get(source, timeout=10)
            resp.raise_for_status()
            return BytesIO(resp.content)
        except Exception as exc:
            warnings.warn(f"Failed to download image from URL '{source}': {exc}")
            return None

    # ── Local file path ─────────────────────────────────────────────────────
    p = Path(source)
    if p.exists() and p.is_file():
        return BytesIO(p.read_bytes())

    # ── Pexels keyword search ────────────────────────────────────────────────
    api_key = os.environ.get("PEXELS_API_KEY")
    if api_key:
        try:
            resp = requests.get(
                "https://api.pexels.com/v1/search",
                headers={"Authorization": api_key},
                params={"query": source, "per_page": 1, "orientation": "landscape"},
                timeout=10,
            )
            resp.raise_for_status()
            photos = resp.json().get("photos", [])
            if not photos:
                warnings.warn(f"Pexels returned no results for query '{source}'.")
                raise ValueError("no results")
            img_url = photos[0]["src"]["large"]
            img_resp = requests.get(img_url, timeout=15)
            img_resp.raise_for_status()
            return BytesIO(img_resp.content)
        except Exception as exc:
            warnings.warn(f"Pexels fetch failed for '{source}': {exc}. Trying picsum fallback.")
    else:
        warnings.warn(
            f"Image source '{source}' is not a URL or local path, and "
            "PEXELS_API_KEY is not set — falling back to picsum."
        )

    # ── Picsum fallback (no key required) ────────────────────────────────────
    try:
        seed = abs(hash(source)) % 1000
        fallback_url = f"https://picsum.photos/seed/{seed}/960/540"
        img_resp = requests.get(fallback_url, timeout=10)
        img_resp.raise_for_status()
        return BytesIO(img_resp.content)
    except Exception as exc2:
        warnings.warn(f"Picsum fallback also failed for '{source}': {exc2}")
        return None

# ─────────────────────────────────────────────
# Image placement helper
# ─────────────────────────────────────────────

def place_image_in_slide(
    slide,
    img_bytes: BytesIO,
    fallback_left: float = 5.0,
    fallback_top:  float = 1.3,
    fallback_width: float = 4.5,
    fallback_height: float = 5.0,
    force_fallback: bool = False,
) -> None:
    """
    Embed `img_bytes` into the slide.
    If `force_fallback` is False, first tries to find a PICTURE or OBJECT
    placeholder in the layout and uses its geometry exactly.
    Pass force_fallback=True to always use the provided geometry (needed when
    the template placeholder covers more area than desired, e.g. half-slide layouts).
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    def _slide_size_inches() -> tuple[float, float]:
        try:
            prs = slide.part.package.presentation_part.presentation
            return prs.slide_width / 914400, prs.slide_height / 914400
        except Exception:
            return 10.0, 5.625

    # Look for a picture / object placeholder in the slide's layout
    pic_ph = None
    if not force_fallback:
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in (
                PP_PLACEHOLDER.PICTURE,
                PP_PLACEHOLDER.OBJECT,
            ):
                pic_ph = ph
                break

    slide_w_in, slide_h_in = _slide_size_inches()

    if pic_ph is not None:
        left, top, width, height = (
            pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
        )
        # Remove the empty placeholder shape before inserting the picture
        sp = pic_ph._element
        sp.getparent().remove(sp)
    else:
        safe_left = max(0.0, min(fallback_left, slide_w_in - 0.5))
        safe_top = max(0.0, min(fallback_top, slide_h_in - 0.5))
        safe_width = max(0.5, min(fallback_width, slide_w_in - safe_left - 0.05))
        safe_height = max(0.5, min(fallback_height, slide_h_in - safe_top - 0.05))

        left = Inches(safe_left)
        top = Inches(safe_top)
        width = Inches(safe_width)
        height = Inches(safe_height)

    # Draw a solid background rect first so any white template-master shapes
    # underneath don't bleed through the chart image.
    _add_filled_rect(
        slide,
        left / 914400, top / 914400,
        width / 914400, height / 914400,
        fill_hex="EFF8F1",
    )

    # ── Crop image to target aspect ratio ('cover' style) ──────────────────
    from PIL import Image as PilImage
    import io
    
    img_bytes.seek(0)
    try:
        img = PilImage.open(img_bytes)
        target_ratio = (width / 914400) / (height / 914400)
        img_w, img_h = img.size
        img_ratio = img_w / img_h
        
        # Center crop the image if aspect ratios don't match
        if abs(target_ratio - img_ratio) > 0.05:
            if img_ratio > target_ratio:
                # Image is too wide: crop left and right
                new_w = int(target_ratio * img_h)
                l_crop = (img_w - new_w) // 2
                img = img.crop((l_crop, 0, l_crop + new_w, img_h))
            else:
                # Image is too tall: crop top and bottom
                new_h = int(img_w / target_ratio)
                t_crop = (img_h - new_h) // 2
                img = img.crop((0, t_crop, img_w, t_crop + new_h))
                
        out = io.BytesIO()
        # Save as PNG to preserve visual quality
        if img.mode in ("RGBA", "P"):
            img.save(out, format="PNG")
        else:
            img.save(out, format="JPEG", quality=90)
        out.seek(0)
        img_bytes = out
    except Exception as e:
        warnings.warn(f"Image cropping failed, using original: {e}")
        img_bytes.seek(0)

    slide.shapes.add_picture(img_bytes, left, top, width, height)


# ─────────────────────────────────────────────
# Drawing primitives
# ─────────────────────────────────────────────

def _add_filled_rect(slide, left_in, top_in, width_in, height_in,
                     fill_hex: str, alpha: float = 1.0) -> None:
    """
    Add a solid-filled rectangle (no border) to the slide.
    `alpha` is approximate — python-pptx does not directly support
    transparency in MSO_FILL, so we skip it silently for values < 1.
    """
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    shape.line.fill.background()  # no border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))


def _add_textbox(
    slide,
    text: str,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    font_size: int = 14,
    bold: bool = False,
    color_hex: str = "#1A1A2E",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    try:
        prs = slide.part.package.presentation_part.presentation
        slide_w_in = prs.slide_width / 914400
        slide_h_in = prs.slide_height / 914400
    except Exception:
        slide_w_in = 10.0
        slide_h_in = 5.625
    safe_left = max(0.0, min(left_in, slide_w_in - 0.4))
    safe_top = max(0.0, min(top_in, slide_h_in - 0.35))
    safe_width = max(0.4, min(width_in, slide_w_in - safe_left - 0.03))
    safe_height = max(0.3, min(height_in, slide_h_in - safe_top - 0.03))

    tx = slide.shapes.add_textbox(
        Inches(safe_left), Inches(safe_top),
        Inches(safe_width), Inches(safe_height),
    )
    tf = tx.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(max(font_size, MIN_TEXT_PT))
    run.font.bold  = bold
    run.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))
    try:
        tf.fit_text(max_size=max(font_size, MIN_TEXT_PT), font_family="Calibri")
    except Exception:
        pass


# ─────────────────────────────────────────────
# Slide-type builders
# ─────────────────────────────────────────────

def add_hero_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Hero / Title slide.

    JSON fields:
        title      (str, required)
        subtitle   (str, optional)
        body       (str, optional)  — e.g. date or tagline
        image      (str, optional)  — URL, local path, or Pexels keyword
        layout_index or layout_name
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    if "subtitle" in data:
        if not fill_placeholder(slide, 1, data["subtitle"], warn_missing=False):
            # Fallback: draw as a textbox below the title area
            _add_textbox(slide, data["subtitle"],
                         0.5, 3.0, 5.5, 1.0,
                         font_size=20, color_hex=TEXT_LIGHT)

    if "body" in data:
        if not fill_placeholder(slide, 2, data["body"], warn_missing=False):
            if not fill_placeholder(slide, 1, data["body"], warn_missing=False):
                _add_textbox(slide, data["body"],
                             0.5, 4.2, 5.5, 0.7,
                             font_size=14, color_hex=TEXT_LIGHT)

    # Optional image — placed in a picture placeholder or right half
    if "image" in data:
        img_bytes = fetch_image_bytes(data["image"])
        if img_bytes:
            place_image_in_slide(
                slide, img_bytes,
                fallback_left=5.5, fallback_top=0.0,
                fallback_width=4.5, fallback_height=7.5,
            )


def add_section_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Section divider slide: just title + subtitle.

    JSON fields:
        title    (str, required)
        subtitle (str, optional)
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    if "subtitle" in data:
        fill_placeholder(slide, 1, data["subtitle"], warn_missing=False)


def add_kpi_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    KPI / Summary slide with coloured metric cards drawn programmatically.

    JSON fields:
        title   (str, required)
        kpis    (list, required)  — list of:
            { "label": str, "value": str, "unit": str, "color": str (optional) }
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", "Highlights"))

    kpis = data.get("kpis", [])
    if not kpis:
        return

    # Grid layout: up to 3 columns, 2 rows
    n     = min(len(kpis), 6)
    n_col = min(n, 3)
    n_row = math.ceil(n / n_col)

    CARD_W    = 2.8
    CARD_H    = 1.8
    START_X   = (10.0 - n_col * CARD_W - (n_col - 1) * 0.2) / 2
    START_Y   = 1.6
    GAP_X     = 0.2
    GAP_Y     = 0.25

    for i, kpi in enumerate(kpis[:n]):
        col = i % n_col
        row = i // n_col

        x = START_X + col * (CARD_W + GAP_X)
        y = START_Y + row * (CARD_H + GAP_Y)
        color = kpi.get("color", BRAND_PALETTE[i % len(BRAND_PALETTE)])

        # Card background
        _add_filled_rect(slide, x, y, CARD_W, CARD_H, fill_hex=color)

        # Value (large centre number)
        _add_textbox(
            slide, kpi.get("value", "—"),
            x + 0.1, y + 0.25, CARD_W - 0.2, 0.8,
            font_size=28, bold=True, color_hex=TEXT_LIGHT,
            align=PP_ALIGN.CENTER,
        )
        # Unit
        _add_textbox(
            slide, kpi.get("unit", ""),
            x + 0.1, y + 0.95, CARD_W - 0.2, 0.35,
            font_size=10, bold=False, color_hex="#D8F3DC",
            align=PP_ALIGN.CENTER,
        )
        # Label
        _add_textbox(
            slide, kpi.get("label", ""),
            x + 0.1, y + 1.3, CARD_W - 0.2, 0.42,
            font_size=11, bold=False, color_hex="#FFFFFF",
            align=PP_ALIGN.CENTER,
        )


def add_chart_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Chart-focused slide with optional text.

    JSON fields:
        title    (str, required)
        subtitle (str, optional)
        body     (str, optional)  — supporting text beside the chart
        bullets  (list, optional) — bullet points beside the chart
        chart    (dict, required) — chart_config passed to render_chart_png()
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    if "subtitle" in data:
        fill_placeholder(slide, 1, data["subtitle"], warn_missing=False)

    chart_cfg = data.get("chart")
    if not chart_cfg:
        warnings.warn(f"Slide id={data.get('id','?')}: type='chart' but no 'chart' key provided.")
        return

    img_buf = render_chart_png(chart_cfg)

    # Attempt to place in a content/object placeholder; otherwise use right zone
    has_left_text = bool(data.get("body") or data.get("bullets"))

    if has_left_text:
        # Chart takes right 60 %, text gets left column
        place_image_in_slide(
            slide, img_buf,
            fallback_left=3.9, fallback_top=1.5,
            fallback_width=5.8, fallback_height=5.5,
        )
        # Text column
        if data.get("bullets"):
            tb_text = "\n".join(f"• {b}" for b in data["bullets"])
        else:
            tb_text = data.get("body", "")

        _add_textbox(
            slide, tb_text,
            0.4, 1.5, 3.3, 5.5,
            font_size=13, color_hex=TEXT_DARK, wrap=True,
        )
    else:
        # Full-width chart
        place_image_in_slide(
            slide, img_buf,
            fallback_left=0.5, fallback_top=1.5,
            fallback_width=9.0, fallback_height=5.5,
        )


def add_image_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Image-focused slide with title and optional body text.

    JSON fields:
        title          (str, required)
        body           (str, optional)
        image          (str, required)  — URL, local path, or Pexels keyword
        image_position (str, optional)  — "right" | "left" | "background" (default: "right")
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    position  = data.get("image_position", "right")
    img_bytes = fetch_image_bytes(data.get("image"))

    if img_bytes:
        if position == "background":
            # Full-bleed background image, then a semi-transparent overlay for readability
            place_image_in_slide(
                slide, img_bytes,
                fallback_left=0.0, fallback_top=0.0,
                fallback_width=10.0, fallback_height=5.625,
                force_fallback=True,
            )
            # Dark overlay strip for text zone
            _add_filled_rect(slide, 0.0, 0.0, 5.2, 5.625, fill_hex="#0D1B2A")

            if data.get("body"):
                _add_textbox(slide, data["body"],
                             0.4, 1.5, 4.6, 3.9, font_size=13, color_hex=TEXT_LIGHT)

        elif position == "left":
            IMG_TOP = 1.3
            img_h = round(5.625 - IMG_TOP - 0.15, 3)
            place_image_in_slide(
                slide, img_bytes,
                fallback_left=0.0, fallback_top=IMG_TOP,
                fallback_width=5.0, fallback_height=img_h,
                force_fallback=True,
            )
            if data.get("body"):
                _add_textbox(slide, data["body"],
                             5.2, IMG_TOP, 4.5, img_h, font_size=13, color_hex=TEXT_DARK)

        else:  # default: right
            IMG_TOP = 1.3
            img_h = round(5.625 - IMG_TOP - 0.15, 3)
            place_image_in_slide(
                slide, img_bytes,
                fallback_left=4.9, fallback_top=IMG_TOP,
                fallback_width=4.9, fallback_height=img_h,
                force_fallback=True,
            )
            if data.get("body"):
                # Skip the white-filled placeholder; always draw a transparent textbox
                _add_textbox(slide, data["body"],
                             0.35, IMG_TOP, 4.35, img_h, font_size=13, color_hex=TEXT_DARK)
    else:
        # No image — fall back to text-only layout
        if data.get("body"):
            _add_textbox(slide, data["body"],
                         0.5, 1.5, 9.0, 3.9, font_size=13, color_hex=TEXT_DARK)


def add_two_column_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Two-column comparison slide.

    Designed for the "TITLE_AND_TWO_COLUMNS" layout (layout_index=3) which has:
        ph_idx=0  TITLE  — slide title
        ph_idx=1  SUBTITLE — left column heading
        ph_idx=2  SUBTITLE — right column heading
        ph_idx=3  SUBTITLE — left column body (bullets)
        ph_idx=4  SUBTITLE — right column body (bullets)

    Falls back to drawing two textboxes if those placeholders are missing.

    JSON fields:
        title (str, required)
        left  (dict, required)  — { "heading": str, "bullets": [str, ...] }
        right (dict, required)  — { "heading": str, "bullets": [str, ...] }
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    left  = data.get("left",  {})
    right = data.get("right", {})

    left_heading  = left.get("heading",  "Left")
    right_heading = right.get("heading", "Right")
    left_bullets  = left.get("bullets",  [])
    right_bullets = right.get("bullets", [])
    left_body     = "\n".join(f"• {b}" for b in left_bullets)
    right_body    = "\n".join(f"• {b}" for b in right_bullets)

    # ── Column headings (ph 1 / ph 2) ────────────────────────────────────────
    if not fill_placeholder(slide, 1, left_heading.upper()):
        _add_textbox(slide, left_heading.upper(), 0.3, 1.4, 4.5, 0.55,
                     font_size=13, bold=True, color_hex=BRAND_PALETTE[1])

    if not fill_placeholder(slide, 2, right_heading.upper()):
        _add_textbox(slide, right_heading.upper(), 5.1, 1.4, 4.5, 0.55,
                     font_size=13, bold=True, color_hex=BRAND_PALETTE[2])

    # ── Column bodies (ph 3 / ph 4) ───────────────────────────────────────────
    # ph 3 and 4 are the content body boxes in the TITLE_AND_TWO_COLUMNS layout.
    # Fallback to explicit textboxes only — never re-use ph 1/ph 2 (headings) as
    # body containers, which would cause heading text to overlap with body text.
    
    # Calculate a safe dynamic font size based on text length to avoid overflows
    left_char_count = sum(len(b) for b in left_bullets)
    right_char_count = sum(len(b) for b in right_bullets)
    max_chars = max(left_char_count, right_char_count)
    dyn_font = 10 if max_chars > 300 else (11 if max_chars > 200 else 12)
    
    if not add_bullets_to_placeholder(slide, 3, left_bullets):
        _add_textbox(slide, left_body, 0.3, 2.0, 4.4, 3.4,
                     font_size=dyn_font, color_hex=TEXT_DARK)

    if not add_bullets_to_placeholder(slide, 4, right_bullets):
        _add_textbox(slide, right_body, 5.1, 2.0, 4.4, 3.4,
                     font_size=dyn_font, color_hex=TEXT_DARK)


def add_bullets_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Standard title + bullets slide.

    JSON fields:
        title   (str, required)
        bullets (list[str], required)
        body    (str, optional)  — alternative to bullets (plain paragraph text)
    """
    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)

    fill_placeholder(slide, 0, data.get("title", ""))

    bullets = data.get("bullets", [])
    body    = data.get("body", "")
    used_ph = False

    if bullets:
        used_ph = add_bullets_to_placeholder(slide, 1, bullets)
    elif body:
        used_ph = fill_placeholder(slide, 1, body, warn_missing=False)

    if not used_ph:
        # Fallback: draw textboxes starting just below title
        BODY_TOP  = 1.25
        BODY_H    = 4.1
        if bullets and len(bullets) >= 5:
            # Split into two visually distinct columns
            mid   = math.ceil(len(bullets) / 2)
            col_w = 4.3
            left_text  = "\n".join(f"\u2022  {b}" for b in bullets[:mid])
            right_text = "\n".join(f"\u2022  {b}" for b in bullets[mid:])
            _add_textbox(slide, left_text,  0.40, BODY_TOP,          col_w, BODY_H,
                         font_size=13, color_hex=TEXT_DARK, wrap=True)
            _add_textbox(slide, right_text, 0.40 + col_w + 0.30, BODY_TOP, col_w, BODY_H,
                         font_size=13, color_hex=TEXT_DARK, wrap=True)
        elif bullets:
            text = "\n".join(f"\u2022  {b}" for b in bullets)
            _add_textbox(slide, text, 0.40, BODY_TOP, 9.2, BODY_H,
                         font_size=14, color_hex=TEXT_DARK, wrap=True)
        elif body:
            _add_textbox(slide, body, 0.40, BODY_TOP, 9.2, BODY_H,
                         font_size=14, color_hex=TEXT_DARK, wrap=True)


# ─────────────────────────────────────────────
# Layout-11 KPI Template slide  (BLANK_1 — 3×2 icon-grid)
# ─────────────────────────────────────────────
# Layout 11 (BLANK_1) placeholder map:
#  Row 1  — value boxes (0.85"×0.63"): ph 2 (x=1.74), ph 4 (x=4.58), ph 6 (x=7.42)
#           label boxes (2.75"×0.58"): ph 1 (x=0.79), ph 8 (x=3.63), ph 9 (x=6.47)
#  Row 2  — value boxes (0.85"×0.58"): ph 3 (x=1.74), ph 5 (x=4.58), ph 7 (x=7.42)
#           label boxes (2.75"×0.58"): ph13 (x=0.79), ph14 (x=3.63), ph15 (x=6.47)

def add_kpi_template_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    KPI summary slide using layout 11 (BLANK_1) which provides a native
    3 × 2 icon-grid.  Each cell in the template has:
      • a small 'value' placeholder (for the big number / percentage)
      • a wider 'label' placeholder (for the metric name + unit)
    The template's own leaf/icon illustrations remain visible underneath.

    JSON fields:
        title  (str, required)
        kpis   (list, required) — up to 6 items:
            { "label": str, "value": str, "unit": str }
    """
    VALUE_SLOTS = [2, 4, 6, 3, 5, 7]      # small boxes — bold number
    LABEL_SLOTS = [1, 8, 9, 13, 14, 15]  # wide boxes  — label + unit

    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)
    fill_placeholder(slide, 0, data.get("title", "Key Metrics"))

    for i, kpi in enumerate(data.get("kpis", [])[:6]):
        v_ph = VALUE_SLOTS[i]
        l_ph = LABEL_SLOTS[i]

        value = kpi.get("value", "—")
        label = kpi.get("label", "")
        unit  = kpi.get("unit",  "")

        # Value — fits in the small decorative-icon overlay box
        fill_placeholder(slide, v_ph, value, warn_missing=False)
        set_placeholder_font(slide, v_ph, bold=True,  size_pt=13,
                             color_hex=BRAND_PALETTE[0])

        # Label + unit on two lines in the wider description box
        label_text = label if not unit else f"{label}\n{unit}"
        fill_placeholder(slide, l_ph, label_text, warn_missing=False)
        set_placeholder_font(slide, l_ph, bold=False, size_pt=10,
                             color_hex=BRAND_PALETTE[1])


# ─────────────────────────────────────────────
# Layout-16 Three-column slide  (BLANK_1_1_1_2)
# ─────────────────────────────────────────────
# ph_idx=0  : title (0.79", 0.49")  8.43"×0.63"
# Column headers (2.56"×0.80"):  ph 4 (x=0.78), ph 5 (x=3.72), ph 6 (x=6.65)
# Column bodies  (2.56"×2.16"):  ph 1 (x=0.78), ph 2 (x=3.72), ph 3 (x=6.65)

def add_three_column_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Three-column analysis/initiative slide using layout 16 (BLANK_1_1_1_2).
    Each column has a header + body area built into the template.

    JSON fields:
        title    (str, required)
        columns  (list, required) — exactly 3 items:
            { "heading": str, "bullets": [str, ...] or "body": str }
    """
    HEADER_SLOTS = [4, 5, 6]
    BODY_SLOTS   = [1, 2, 3]

    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)
    fill_placeholder(slide, 0, data.get("title", ""))

    columns = data.get("columns", [])
    for i, col in enumerate(columns[:3]):
        h_ph = HEADER_SLOTS[i]
        b_ph = BODY_SLOTS[i]

        # Column header
        heading = col.get("heading", "")
        fill_placeholder(slide, h_ph, heading.upper(), warn_missing=False)
        set_placeholder_font(slide, h_ph, bold=True, size_pt=11,
                             color_hex=BRAND_PALETTE[0])

        # Column body — bullets or plain text
        col_bullets = col.get("bullets", [])
        col_body    = col.get("body", "")
        if col_bullets:
            body_text = "\n".join(f"\u2022 {b}" for b in col_bullets)
        else:
            body_text = col_body

        if not fill_placeholder(slide, b_ph, body_text, warn_missing=False):
            # Fallback textbox at the column's approximate x-position
            col_x = [0.78, 3.72, 6.65][i]
            _add_textbox(slide, body_text, col_x, 2.75, 2.56, 2.20,
                         font_size=11, color_hex=TEXT_DARK, wrap=True)
        else:
            set_placeholder_font(slide, b_ph, bold=False, size_pt=11,
                                 color_hex=TEXT_DARK)


# ─────────────────────────────────────────────
# Layout-17 Steps / Cascade slide  (BLANK_1_1_1_2_1)
# ─────────────────────────────────────────────
# ph_idx=0  : title (0.79", 0.49")  8.43"×0.63"
# Step headings (5.84"×0.50"): ph 4 (y=1.24), ph 5 (y=2.58), ph 6 (y=3.91)
# Step bodies   (5.84"×0.58"): ph 2 (y=1.73), ph 3 (y=3.06), ph 1 (y=4.40)
# x-offsets cascade: 2.14", 2.55", 2.97"  — template shows staircase arrows

def add_steps_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """
    Cascade / step-progression slide using layout 17 (BLANK_1_1_1_2_1).
    The template renders three staircase-style numbered steps.

    JSON fields:
        title  (str, required)
        steps  (list, required) — up to 3 items:
            { "heading": str, "body": str }
    """
    HEADING_SLOTS = [4, 5, 6]
    BODY_SLOTS    = [2, 3, 1]

    layout = get_layout(prs, data)
    slide  = prs.slides.add_slide(layout)
    fill_placeholder(slide, 0, data.get("title", ""))

    for i, step in enumerate(data.get("steps", [])[:3]):
        h_ph = HEADING_SLOTS[i]
        b_ph = BODY_SLOTS[i]

        heading = step.get("heading", f"Step {i + 1}")
        body    = step.get("body", "")

        fill_placeholder(slide, h_ph, heading.upper(), warn_missing=False)
        set_placeholder_font(slide, h_ph, bold=True, size_pt=12,
                             color_hex=BRAND_PALETTE[0])

        fill_placeholder(slide, b_ph, body, warn_missing=False)
        set_placeholder_font(slide, b_ph, bold=False, size_pt=11,
                             color_hex=TEXT_DARK)


# ─────────────────────────────────────────────
# Dispatcher
# ─────────────────────────────────────────────

SLIDE_BUILDERS = {
    "hero":           add_hero_slide,
    "section":        add_section_slide,
    "kpi_summary":    add_kpi_slide,
    "kpi_template":   add_kpi_template_slide,  # layout 11 — native 3×2 icon grid
    "chart":          add_chart_slide,
    "image":          add_image_slide,
    "two_column":     add_two_column_slide,
    "three_column":   add_three_column_slide,  # layout 16 — native 3-col
    "steps":          add_steps_slide,          # layout 17 — cascade steps
    "bullets":        add_bullets_slide,
}


def generate_presentation(
    plan: dict[str, Any],
    template_path: str | None = None,
    output_path:   str | None = None,
) -> str:
    """
    Main orchestrator.

    Loads the template, iterates the slide plan, calls the appropriate
    builder for each slide type, and saves the output .pptx.

    Returns the resolved output path string.
    """
    meta     = plan.get("meta", {})
    tmpl     = template_path or meta.get("template", "template.pptx")
    out      = output_path   or meta.get("output",   "output.pptx")
    slides   = plan.get("slides", [])

    if not Path(tmpl).exists():
        raise FileNotFoundError(
            f"Template not found: '{tmpl}'\n"
            "Pass --template <path> or set 'meta.template' in your JSON plan."
        )

    print(f"  Template : {tmpl}")
    print(f"  Output   : {out}")
    print(f"  Slides   : {len(slides)}")
    print()

    # Load the template as the base presentation, then strip its sample slides.
    # The slide masters, layouts, colour theme, and fonts are all preserved —
    # only the pre-existing content slides are removed so we start with a clean deck.
    prs = Presentation(tmpl)
    _clear_template_slides(prs)

    for slide_data in slides:
        slide_id   = slide_data.get("id", "?")
        slide_type = slide_data.get("type", "bullets")
        builder    = SLIDE_BUILDERS.get(slide_type)

        if builder is None:
            warnings.warn(
                f"  [slide {slide_id}] Unknown type='{slide_type}' — "
                f"supported: {', '.join(SLIDE_BUILDERS)}. Skipping."
            )
            continue

        try:
            print(f"  [{slide_id:>2}]  type={slide_type:<14}  "
                  f"layout={slide_data.get('layout_index', slide_data.get('layout_name', 'auto'))}")
            builder(prs, slide_data)
        except Exception as exc:
            warnings.warn(f"  [slide {slide_id}] Error building slide: {exc}")
            import traceback
            traceback.print_exc()

    prs.save(out)
    print(f"\n  Saved → {out}")
    return out


# ─────────────────────────────────────────────
# CLI entry point
# ─────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Template-based PPTX generator using python-pptx + Matplotlib/Seaborn.\n\n"
            "Examples:\n"
            "  python generate_from_template.py --inspect-template \"template.pptx\"\n"
            "  python generate_from_template.py plan.json\n"
            "  python generate_from_template.py plan.json --output report.pptx\n"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )

    parser.add_argument(
        "plan",
        nargs="?",
        help="Path to the slide plan JSON file.",
    )
    parser.add_argument(
        "--inspect-template",
        metavar="TEMPLATE",
        help="Print all slide layouts and placeholders in the given template file, then exit.",
    )
    parser.add_argument(
        "--template",
        metavar="PATH",
        help="Override the template path (default: value from plan.meta.template).",
    )
    parser.add_argument(
        "--output",
        metavar="PATH",
        help="Override the output .pptx path (default: value from plan.meta.output).",
    )

    args = parser.parse_args()

    # ── Inspect mode ────────────────────────────────────────────────────────
    if args.inspect_template:
        tmpl_path = args.inspect_template
        if not Path(tmpl_path).exists():
            print(f"ERROR: Template file not found: '{tmpl_path}'", file=sys.stderr)
            sys.exit(1)
        prs = Presentation(tmpl_path)
        list_layouts(prs)
        sys.exit(0)

    # ── Generate mode ───────────────────────────────────────────────────────
    if not args.plan:
        parser.print_help()
        sys.exit(1)

    plan_path = Path(args.plan)
    if not plan_path.exists():
        print(f"ERROR: Slide plan not found: '{plan_path}'", file=sys.stderr)
        sys.exit(1)

    with open(plan_path, encoding="utf-8") as f:
        plan = json.load(f)

    print(f"\n{'─'*50}")
    print("  ESG PPTX Generator — from template")
    print(f"{'─'*50}")

    generate_presentation(plan, args.template, args.output)


if __name__ == "__main__":
    main()
