"""
generate_ppt.py — Python/FastAPI equivalent of app/api/generate-ppt/route.ts
Converts PptxGenJS slide-building to python-pptx.
Dependencies: fastapi uvicorn python-pptx httpx lxml pydantic
"""

# ── Standard library ──────────────────────────────────────────────────────────
import asyncio
import base64
import io
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Optional
from urllib.parse import quote_plus


# Load .env.local (Next.js convention) then .env from the project root upward
def _load_env() -> None:
    root = Path(__file__).resolve()
    for directory in [root.parent, *root.parents]:
        for filename in (".env.local", ".env"):
            p = directory / filename
            if p.exists():
                try:
                    from dotenv import load_dotenv

                    load_dotenv(p, override=False)
                    print(f"Loaded env from {p}")
                except ImportError:
                    # Manual fallback if python-dotenv not installed
                    for line in p.read_text(encoding="utf-8").splitlines():
                        line = line.strip()
                        if not line or line.startswith("#") or "=" not in line:
                            continue
                        k, _, v = line.partition("=")
                        k = k.strip()
                        v = v.strip().strip('"').strip("'")
                        if k and k not in os.environ:
                            os.environ[k] = v
                    print(f"Loaded env from {p} (manual)")
                return


_load_env()

# ── Third-party ───────────────────────────────────────────────────────────────
import httpx
from fastapi import FastAPI
from fastapi.responses import Response
from pydantic import BaseModel

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from lxml import etree

# ── DrawingML XML namespace ───────────────────────────────────────────────────
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# OPC relationship namespace (used for slide removal)
_R_NS_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _clear_template_slides(prs: Presentation) -> None:
    """
    Remove every sample slide from the loaded template while keeping
    slide masters, layouts, theme, and fonts intact.
    Properly drops both the OPC relationship AND the sldId XML element
    so no orphaned parts remain in the output package.
    """
    sldIdLst = prs.slides._sldIdLst
    rIds = [sld.get(f"{{{_R_NS_REL}}}id") for sld in list(sldIdLst)]
    for rId in rIds:
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)


# ═══════════════════════════════════════════════════════════════════════════════
# Brand constants  (mirrors TS constants at the top of route.ts)
# ═══════════════════════════════════════════════════════════════════════════════
BRAND_GREEN = "3DAF5C"  # mid-green  — borders, accents, bullet dots
BRAND_BLUE = "2A7FD5"
HEADER_BG = "1E3A4F"  # legacy (kept for route.ts compat)
BODY_BG = "FFFFFF"
TEXT_DARK = "1A1A2E"
TEXT_MID = "555566"

# ── Slidesgo "Sustainability & Environment" theme palette ─────────────────────
# Matches the light sage-green template shown in the design brief.
HEADING_GREEN = "1E3A0E"  # dark forest green  — all slide headings & labels
ACCENT_ORANGE = "F5A623"  # warm amber/orange  — icons, KPI values, highlights
BG_LIGHT = "D8ECC8"  # light sage green   — slide body background
CONTENT_BOX = "EEF7E4"  # off-white green    — card & table-cell fills
BODY_TEXT = "3D4A2E"  # dark olive         — body paragraphs & bullets

W = 10.0  # slide width  (inches, 16:9)
H = 5.625  # slide height (inches, 16:9)

LOGO_PATH = (
    Path(__file__).resolve().parent.parent.parent.parent
    / "Growlity-Logo-Website-2-new.png.webp"
)

# Slidesgo template — loaded once at startup so all generated decks inherit the
# template's slide master, colour theme, fonts, and decorative design elements.
_TEMPLATE_PATH = (
    Path(__file__).resolve().parent.parent.parent.parent
    / "Sustainability and Environment Newsletter _ by Slidesgo.pptx"
)
_TEMPLATE_BLANK_LAYOUT_IDX = 10  # "BLANK" – no placeholders
_TEMPLATE_TITLE_LAYOUT_IDX = (
    0  # "TITLE" – CENTER_TITLE (ph0) + SUBTITLE (ph1), used for hero
)
_TEMPLATE_CONTENT_LAYOUT_IDX = (
    4  # "TITLE_ONLY" – TITLE (ph0), used for all content slides
)


def _read_template_fonts() -> tuple[str, str]:
    """Extract heading and body font names from the Slidesgo template's theme XML."""
    try:
        import zipfile

        with zipfile.ZipFile(str(_TEMPLATE_PATH), "r") as z:
            xml_bytes = z.read("ppt/theme/theme1.xml")
        root = etree.fromstring(xml_bytes)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        major = root.find(".//a:majorFont/a:latin", ns)
        minor = root.find(".//a:minorFont/a:latin", ns)
        major_face = major.get("typeface", "Lato") if major is not None else "Lato"
        minor_face = minor.get("typeface", "Lato") if minor is not None else "Lato"
        # Ignore generic theme references like "+mj-lt" / "+mn-lt"
        if major_face.startswith("+") or not major_face:
            major_face = "Lato"
        if minor_face.startswith("+") or not minor_face:
            minor_face = "Lato"
        print(f"Template fonts: heading={major_face!r}, body={minor_face!r}")
        return major_face, minor_face
    except Exception as exc:
        print(f"Could not read template fonts: {exc} — defaulting to Lato")
        return "Lato", "Lato"


if _TEMPLATE_PATH.exists():
    TEMPLATE_FONT_MAJOR, TEMPLATE_FONT_BODY = _read_template_fonts()
else:
    TEMPLATE_FONT_MAJOR, TEMPLATE_FONT_BODY = "Lato", "Lato"

PEXELS_API_KEY: str = os.environ.get("PEXELS_API_KEY", "")
CHUTES_API_KEY: str = (
    os.environ.get("CHUTES_API_KEY", "")
    or os.environ.get("NEXT_PUBLIC_CHUTES_API_KEY", "")
    or os.environ.get("OPENAI_API_KEY", "")  # fallback: key name used in env.txt
)
CHUTES_URL = "https://llm.chutes.ai/v1/chat/completions"

# Color palettes — sustainability / nature theme to complement the Slidesgo template
# Palette: dark forest green, warm orange, mid/light greens — matches template screenshots
PALETTES = [
    [
        "#1E3A0E",
        "#F5A623",
        "#3DAF5C",
        "#8DB87A",
        "#F08C00",
        "#A3C97A",
    ],  # forest + orange
    ["#2E7D52", "#F5A623", "#1A8C6A", "#3DAF7A", "#F08C00", "#4ABC8A"],  # teal + orange
    [
        "#2E7D32",
        "#F5A623",
        "#43A047",
        "#66BB6A",
        "#E65100",
        "#A5D6A7",
    ],  # greens + orange
    ["#5B8C3A", "#F5A623", "#7AB545", "#4A7D2A", "#E67E00", "#8ABF5A"],  # sage + orange
    [
        "#1F6B5A",
        "#F5A623",
        "#2E8C6A",
        "#3D9E7A",
        "#F08C00",
        "#4AAB8A",
    ],  # deep teal + orange
]
PALETTE_BG = ["#EEF7E4", "#EDF7ED", "#EBF5ED", "#EEF7E4", "#EBF5F0"]

MIN_BODY_FONT = 10
MIN_BULLET_FONT = 9.5
MAX_BULLETS_SINGLE_COL = 5
MAX_BULLETS_TWO_COL = 8
MAX_BULLETS_PER_COLUMN = 4  # Maximum bullets per column to prevent overflow


# ═══════════════════════════════════════════════════════════════════════════════
# Low-level helpers
# ═══════════════════════════════════════════════════════════════════════════════


def hex_color(h: str) -> RGBColor:
    """Convert a 6-char hex string (no '#') to RGBColor."""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _alpha_hex(hex_color_str: str, pct: float) -> str:
    """Append an 8-bit alpha suffix to a hex color string (for QuickChart)."""
    h = hex_color_str.lstrip("#")
    alpha = format(round(pct * 255), "02x")
    return f"#{h}{alpha}"


def add_rect(slide, x: float, y: float, w: float, h: float, color_hex: str):
    """
    Add a solid-filled rectangle with no visible border.
    Replaces slide.addShape("rect", {x, y, w, h, fill, line}) from PptxGenJS.
    All coordinates in inches.
    """
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    # Solid fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(color_hex)
    # Remove border by injecting <a:noFill> into the line element via XML
    sp_pr = shape._element.spPr
    ln = sp_pr.get_or_add_ln()
    # Clear existing children and insert noFill
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, f"{{{_A_NS}}}noFill")
    return shape


def _add_slide_dot_clusters(slide) -> None:
    """
    Draw two 3×3 grids of small BRAND_GREEN squares — the dot-grid (.  .  .  .)
    motif visible in the Slidesgo template corners.
    Top-left cluster at (0.05", 0.05") and bottom-left cluster above the footer.
    Called after the BG_LIGHT background rect is placed so the dots stay on top.
    """
    dot_size = 0.055  # square side length (inches)
    dot_gap = 0.082  # center-to-center spacing (inches)
    # Top-left cluster
    for row in range(3):
        for col in range(3):
            add_rect(
                slide,
                0.05 + col * dot_gap,
                0.05 + row * dot_gap,
                dot_size,
                dot_size,
                BRAND_GREEN,
            )
    # Bottom-left cluster (above the footer / bottom edge)
    bl_y = H - 3 * dot_gap - 0.07
    for row in range(3):
        for col in range(3):
            add_rect(
                slide,
                0.05 + col * dot_gap,
                bl_y + row * dot_gap,
                dot_size,
                dot_size,
                BRAND_GREEN,
            )


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 11,
    bold: bool = False,
    color_hex: str = TEXT_DARK,
    font_face: str = "",  # empty → resolved to TEMPLATE_FONT_BODY at call time
    align: str = "left",
    italic: bool = False,
    valign: str = "top",
    wrap: bool = True,
):
    """
    Add a text box with a single run of styled text.
    Replaces slide.addText(text, {x, y, w, h, ...}) from PptxGenJS.
    All coordinates in inches.
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # Vertical alignment
    _valign_map = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    tf.vertical_anchor = _valign_map.get(valign, MSO_ANCHOR.TOP)

    # Text and run formatting on the first paragraph
    para = tf.paragraphs[0]
    _align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    para.alignment = _align_map.get(align, PP_ALIGN.LEFT)

    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_face or TEMPLATE_FONT_BODY
    run.font.color.rgb = hex_color(color_hex)

    return txBox


def add_picture_from_bytes(
    slide, img_bytes: bytes, x: float, y: float, w: float, h: float
):
    """
    Insert an image from raw bytes with "cover" style sizing (center crop to match target aspect ratio).
    Replaces slide.addImage({data: "mime;base64,...", x, y, w, h}) from PptxGenJS.
    All coordinates in inches.
    """
    from PIL import Image as PilImage

    buf = io.BytesIO(img_bytes)
    buf.seek(0)
    try:
        img = PilImage.open(buf)
        fmt = img.format or ""
        
        # Target aspect ratio
        target_ratio = w / h
        img_w, img_h = img.size
        img_ratio = img_w / img_h
        
        # Center crop the image if aspect ratios don't match
        if abs(target_ratio - img_ratio) > 0.05:
            if img_ratio > target_ratio:
                # Image is too wide: crop left and right
                new_w = int(target_ratio * img_h)
                left = (img_w - new_w) // 2
                img = img.crop((left, 0, left + new_w, img_h))
            else:
                # Image is too tall: crop top and bottom
                new_h = int(img_w / target_ratio)
                top = (img_h - new_h) // 2
                img = img.crop((0, top, img_w, top + new_h))
                
        # Convert unsupported formats
        if fmt.upper() not in ("BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"):
            out = io.BytesIO()
            img.convert("RGBA").save(out, format="PNG")
            out.seek(0)
            buf = out
        else:
            # We crop/modified it, but it might be supported natively. 
            # Easiest to just save as PNG again if cropped.
            if hasattr(img, "is_animated") and img.is_animated:
                # Can't crop GIF easily if animated, just pass Original Buffer
                buf.seek(0)
            else:
                out = io.BytesIO()
                img.save(out, format="PNG")
                out.seek(0)
                buf = out
    except Exception:
        buf.seek(0)
    return slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))


def render_icon_png(
    icon_name: str, size_px: int = 64, color_hex: str = "1E3A0E"
) -> Optional[bytes]:
    """
    Render a simple ESG-themed outline icon as a transparent-background PNG.
    Draws at 4× supersampling then downsamples with LANCZOS for clean anti-aliasing.
    Returns PNG bytes, or None on error.
    Supported icons: factory, globe, bolt, leaf, recycle, shield, award,
                     trending-up, users, building, cloud, sun, water, flame, chart-bar, tree
    """
    try:
        import math
        from PIL import Image, ImageDraw

        h_s = color_hex.lstrip("#")
        r, g, b = int(h_s[0:2], 16), int(h_s[2:4], 16), int(h_s[4:6], 16)
        fg = (r, g, b, 255)

        SC = 4  # supersampling scale factor
        S = size_px * SC  # canvas size in pixels
        pad = S * 0.10  # padding
        D = S - 2 * pad  # drawable width/height within padding
        lw = max(SC * 2, S // 13)  # stroke width, scales with icon size

        img = Image.new("RGBA", (S, S), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)

        def px(nx: float, ny: float) -> tuple:
            """Normalised 0-1 coordinate → pixel coordinate within padded canvas."""
            return (pad + nx * D, pad + ny * D)

        def arc_pts(cx, cy, rx, ry, a0, a1, n=40):
            pts = []
            for i in range(n + 1):
                t = math.radians(a0 + (a1 - a0) * i / n)
                pts.append((cx + rx * math.cos(t), cy + ry * math.sin(t)))
            return pts

        if icon_name == "factory":
            draw.rectangle([px(0.05, 0.48), px(0.95, 0.95)], outline=fg, width=lw)
            draw.rectangle(
                [px(0.10, 0.18), px(0.26, 0.48)], outline=fg, width=lw
            )  # smokestack L
            draw.rectangle(
                [px(0.36, 0.22), px(0.52, 0.48)], outline=fg, width=lw
            )  # smokestack R
            draw.rectangle(
                [px(0.12, 0.56), px(0.34, 0.75)], outline=fg, width=lw
            )  # window L
            draw.rectangle(
                [px(0.44, 0.56), px(0.64, 0.75)], outline=fg, width=lw
            )  # window C
            draw.rectangle(
                [px(0.72, 0.60), px(0.88, 0.95)], outline=fg, width=lw
            )  # door

        elif icon_name == "globe":
            draw.ellipse([px(0.02, 0.02), px(0.98, 0.98)], outline=fg, width=lw)
            draw.line([px(0.02, 0.50), px(0.98, 0.50)], fill=fg, width=lw)  # equator
            draw.line(
                [px(0.50, 0.02), px(0.50, 0.98)], fill=fg, width=lw
            )  # prime meridian
            cx, cy = pad + D * 0.50, pad + D * 0.28
            draw.line(
                arc_pts(cx, cy, D * 0.48, D * 0.16, 0, 180), fill=fg, width=lw
            )  # tropic N
            cx2, cy2 = pad + D * 0.50, pad + D * 0.72
            draw.line(
                arc_pts(cx2, cy2, D * 0.48, D * 0.16, 0, 180), fill=fg, width=lw
            )  # tropic S

        elif icon_name == "bolt":
            pts = [
                px(0.62, 0.0),
                px(0.35, 0.46),
                px(0.54, 0.46),
                px(0.38, 1.0),
                px(0.65, 0.54),
                px(0.46, 0.54),
            ]
            draw.polygon(pts, outline=fg, width=lw)

        elif icon_name == "leaf":
            cx_l = pad + D * 0.5
            cy_l = pad + D * 0.5
            draw.line(
                arc_pts(cx_l, cy_l, D * 0.44, D * 0.44, -110, 70), fill=fg, width=lw
            )
            draw.line(
                arc_pts(cx_l, cy_l, D * 0.44, D * 0.44, 110, 290), fill=fg, width=lw
            )
            draw.line([px(0.50, 0.06), px(0.50, 0.94)], fill=fg, width=lw)  # midrib

        elif icon_name == "recycle":
            cx_r = pad + D * 0.5
            cy_r = pad + D * 0.5
            for i, a0a1 in enumerate([(30, 130), (150, 250), (270, 370)]):
                pts = arc_pts(cx_r, cy_r, D * 0.40, D * 0.40, a0a1[0], a0a1[1], n=20)
                draw.line(pts, fill=fg, width=lw)
                end = pts[-1]
                ha = math.radians(a0a1[1] + 90)
                h1 = (
                    end[0] + D * 0.07 * math.cos(ha - 0.5),
                    end[1] + D * 0.07 * math.sin(ha - 0.5),
                )
                h2 = (
                    end[0] + D * 0.07 * math.cos(ha + 0.5),
                    end[1] + D * 0.07 * math.sin(ha + 0.5),
                )
                draw.polygon([end, h1, h2], fill=fg)

        elif icon_name == "shield":
            top_pts = arc_pts(
                pad + D * 0.5, pad + D * 0.42, D * 0.46, D * 0.40, 180, 360, n=30
            )
            all_pts = top_pts + [px(0.96, 0.56), px(0.50, 0.98), px(0.04, 0.56)]
            draw.line(all_pts + [all_pts[0]], fill=fg, width=lw)

        elif icon_name == "award":
            draw.ellipse([px(0.20, 0.02), px(0.80, 0.60)], outline=fg, width=lw)
            draw.line(
                [
                    px(0.26, 0.58),
                    px(0.12, 0.98),
                    px(0.50, 0.76),
                    px(0.88, 0.98),
                    px(0.74, 0.58),
                ],
                fill=fg,
                width=lw,
            )

        elif icon_name == "trending-up":
            draw.line(
                [px(0.05, 0.86), px(0.35, 0.62), px(0.62, 0.42), px(0.90, 0.16)],
                fill=fg,
                width=lw,
            )
            draw.line(
                [px(0.90, 0.16), px(0.70, 0.16)], fill=fg, width=lw
            )  # arrowhead H
            draw.line(
                [px(0.90, 0.16), px(0.90, 0.36)], fill=fg, width=lw
            )  # arrowhead V
            draw.line([px(0.04, 0.92), px(0.96, 0.92)], fill=fg, width=lw)  # x-axis
            draw.line([px(0.04, 0.92), px(0.04, 0.08)], fill=fg, width=lw)  # y-axis

        elif icon_name == "users":
            draw.ellipse(
                [px(0.08, 0.04), px(0.42, 0.40)], outline=fg, width=lw
            )  # head L
            draw.line(
                arc_pts(
                    pad + D * 0.25, pad + D * 0.90, D * 0.28, D * 0.28, 180, 360, n=20
                ),
                fill=fg,
                width=lw,
            )  # body L arc
            draw.ellipse(
                [px(0.58, 0.04), px(0.92, 0.40)], outline=fg, width=lw
            )  # head R
            draw.line(
                arc_pts(
                    pad + D * 0.75, pad + D * 0.90, D * 0.28, D * 0.28, 180, 360, n=20
                ),
                fill=fg,
                width=lw,
            )  # body R arc

        elif icon_name == "building":
            draw.rectangle([px(0.10, 0.05), px(0.90, 0.95)], outline=fg, width=lw)
            for wr in [0.14, 0.32, 0.50, 0.68]:
                for wc in [0.17, 0.42, 0.66]:
                    draw.rectangle(
                        [px(wc, wr), px(wc + 0.18, wr + 0.12)], outline=fg, width=lw
                    )
            draw.rectangle(
                [px(0.38, 0.75), px(0.62, 0.95)], outline=fg, width=lw
            )  # door

        elif icon_name == "cloud":
            draw.ellipse([px(0.08, 0.28), px(0.50, 0.80)], outline=fg, width=lw)
            draw.ellipse([px(0.32, 0.16), px(0.74, 0.68)], outline=fg, width=lw)
            draw.ellipse([px(0.54, 0.32), px(0.92, 0.80)], outline=fg, width=lw)
            draw.line([px(0.08, 0.76), px(0.92, 0.76)], fill=fg, width=lw)  # flat base

        elif icon_name == "sun":
            draw.ellipse([px(0.28, 0.28), px(0.72, 0.72)], outline=fg, width=lw)
            for ang in range(0, 360, 45):
                a = math.radians(ang)
                x1 = pad + D * (0.5 + 0.27 * math.cos(a))
                y1 = pad + D * (0.5 + 0.27 * math.sin(a))
                x2 = pad + D * (0.5 + 0.44 * math.cos(a))
                y2 = pad + D * (0.5 + 0.44 * math.sin(a))
                draw.line([(x1, y1), (x2, y2)], fill=fg, width=lw)

        elif icon_name == "water":
            for ri, y_n in enumerate([0.24, 0.50, 0.74]):
                pts = []
                for i in range(40):
                    x = i / 39
                    y = y_n + 0.07 * math.sin(x * math.pi * 3 + ri * 0.8)
                    pts.append(px(x, y))
                draw.line(pts, fill=fg, width=lw)

        elif icon_name == "flame":
            outer = [
                px(0.50, 0.02),
                px(0.74, 0.28),
                px(0.82, 0.17),
                px(0.84, 0.50),
                px(0.70, 0.74),
                px(0.50, 0.98),
                px(0.30, 0.74),
                px(0.16, 0.50),
                px(0.18, 0.17),
                px(0.26, 0.28),
            ]
            draw.polygon(outer, outline=fg, width=lw)
            inner = [
                px(0.50, 0.26),
                px(0.62, 0.44),
                px(0.60, 0.70),
                px(0.50, 0.80),
                px(0.40, 0.70),
                px(0.38, 0.44),
            ]
            draw.line(inner + [inner[0]], fill=fg, width=lw)

        elif icon_name == "chart-bar":
            draw.rectangle([px(0.08, 0.55), px(0.34, 0.92)], outline=fg, width=lw)
            draw.rectangle([px(0.40, 0.30), px(0.66, 0.92)], outline=fg, width=lw)
            draw.rectangle([px(0.72, 0.08), px(0.92, 0.92)], outline=fg, width=lw)
            draw.line([px(0.04, 0.92), px(0.96, 0.92)], fill=fg, width=lw)  # baseline

        elif icon_name == "tree":
            draw.polygon(
                [px(0.50, 0.02), px(0.18, 0.38), px(0.82, 0.38)], outline=fg, width=lw
            )
            draw.polygon(
                [px(0.50, 0.22), px(0.12, 0.60), px(0.88, 0.60)], outline=fg, width=lw
            )
            draw.polygon(
                [px(0.50, 0.44), px(0.08, 0.82), px(0.92, 0.82)], outline=fg, width=lw
            )
            draw.rectangle(
                [px(0.42, 0.82), px(0.58, 0.98)], outline=fg, width=lw
            )  # trunk

        else:
            # Fallback: circle with checkmark
            draw.ellipse([px(0.05, 0.05), px(0.95, 0.95)], outline=fg, width=lw)
            draw.line(
                [px(0.25, 0.5), px(0.45, 0.72), px(0.75, 0.28)], fill=fg, width=lw
            )

        # Downsample to target size with LANCZOS for anti-aliasing
        result = img.resize((size_px, size_px), Image.LANCZOS)
        buf = io.BytesIO()
        result.save(buf, format="PNG")
        buf.seek(0)
        return buf.read()
    except Exception as exc:
        print(f"render_icon_png error ({icon_name}): {exc}")
        return None


def add_bullet_textbox(
    slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = MIN_BULLET_FONT,
    color_hex: str = TEXT_DARK,
    font_face: str = "",  # empty → resolved to TEMPLATE_FONT_BODY at call time
    bullet_color_hex: str = BRAND_GREEN,
    line_spacing: float = 1.35,
    para_space_after: float = 6,
):
    """
    Add a text box containing bullet-point paragraphs.
    Bullet characters are injected via lxml XML because python-pptx has no
    native bullet API.
    Replaces slide.addText(makeBulletItems(bArr), {...}) from PptxGenJS.
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    safe_font_size = max(font_size, MIN_BULLET_FONT)

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.space_after = Pt(para_space_after)
        para.line_spacing = line_spacing

        # Run with text
        run = para.add_run()
        run.text = bullet_text
        run.font.size = Pt(safe_font_size)
        run.font.name = font_face or TEMPLATE_FONT_BODY
        run.font.color.rgb = hex_color(color_hex)

        # Inject bullet character and bullet color via DrawingML XML
        pPr = para._p.get_or_add_pPr()

        # <a:buClr><a:srgbClr val="3DAF5C"/></a:buClr>
        buClr = etree.SubElement(pPr, f"{{{_A_NS}}}buClr")
        srgbClr = etree.SubElement(buClr, f"{{{_A_NS}}}srgbClr")
        srgbClr.set("val", bullet_color_hex.lstrip("#"))

        # <a:buChar char="•"/>
        buChar = etree.SubElement(pPr, f"{{{_A_NS}}}buChar")
        buChar.set("char", "•")

        # Indent so text doesn't overlap the bullet
        pPr.set("indent", str(int(Inches(0.15))))  # hanging indent
        pPr.set("marL", str(int(Inches(0.25))))  # left margin

    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# Text / tag parsing helpers (direct ports of the TS functions)
# ═══════════════════════════════════════════════════════════════════════════════


def strip_md(text: str) -> str:
    """Strip common Markdown syntax that renders literally in presentations."""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"__([^_]+)__", r"\1", text)
    text = re.sub(r"_([^_]+)_", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^[\d]+\.\s*", "", text)
    return text.strip()


def _clean_bullet_prefix(text: str) -> str:
    """Remove AI-generated label prefixes like 'ACTION:', 'INSIGHT:', etc."""
    return re.sub(
        r"^(?:ACTION|INSIGHT|RECOMMENDATION|KEY\s*POINT|NOTE|FINDING|STRATEGY|RISK|TAKEAWAY|PRIORITY|HIGHLIGHT|SUMMARY|IMPLICATION)\s*:\s*",
        "",
        text,
        flags=re.IGNORECASE,
    ).strip()


def _reformat_stat_bullets(bullets: list[str]) -> list[str]:
    """
    Convert STAT-format bullets used in non-summary slides to readable form.
    'STAT: Label | value | unit'  → 'Label: value (unit)'
    'Label | value | unit'        → 'Label: value (unit)'
    Plain bullets are returned unchanged.
    """
    out = []
    for b in bullets:
        m = re.match(r"^STAT:\s*([^|]+)\|\s*([^|]+)\|\s*(.+)$", b, re.IGNORECASE)
        if m:
            out.append(
                f"{m.group(1).strip()}: {m.group(2).strip()} ({m.group(3).strip()})"
            )
            continue
        parts = [p.strip() for p in b.split("|")]
        if len(parts) == 3:
            out.append(f"{parts[0]}: {parts[1]} ({parts[2]})")
        elif len(parts) == 2:
            out.append(f"{parts[0]}: {parts[1]}")
        else:
            out.append(b)
    return out


def parse_chart_tag(text: str) -> Optional[dict]:
    """
    Parse [CHART:type:Label1=val1,Label2=val2] from slide markdown.
    Returns {"type": str, "labels": list[str], "values": list[float]} or None.
    """
    match = re.search(
        r"\[CHART:(pie|doughnut|bar|col|line|area|radar|polar|stacked):([^\]]+)\]",
        text,
        re.IGNORECASE,
    )
    if not match:
        return None
    chart_type = match.group(1).lower()
    pairs = [p.strip() for p in match.group(2).split(",")]
    labels, values = [], []
    for pair in pairs:
        parts = pair.split("=")
        if len(parts) == 2:
            label, val = parts[0].strip(), parts[1].strip()
            try:
                values.append(float(val))
                labels.append(label)
            except ValueError:
                pass
    return {"type": chart_type, "labels": labels, "values": values} if labels else None


def parse_image_tag(text: str) -> Optional[dict]:
    """
    Parse [IMAGE:search query:layout] from slide markdown.
    Returns {"query": str, "layout": str} or None.
    """
    match = re.search(
        r"\[IMAGE:([^\]:]+?)(?::(bg|right|left|inset))?\]", text, re.IGNORECASE
    )
    if not match:
        return None
    return {
        "query": match.group(1).strip(),
        "layout": (match.group(2) or "right").lower(),
    }


def _is_dsl_tag(line: str) -> bool:
    """Return True for lines that are DSL control tags (not user-visible text)."""
    s = line.strip()
    return bool(
        re.match(r"^\[(?:TYPE|IMAGE|CHART):", s, re.IGNORECASE)
        or re.match(r"^\[(?:A|B)\]", s, re.IGNORECASE)
        or re.match(
            r"^\[SLIDE\b", s, re.IGNORECASE
        )  # filter [SLIDE N] slide-number markers
    )


def _split_dense_slide_markdown(slide_md: str) -> list[str]:
    """
    Split very dense text-only slides into two continuation slides to preserve readability.
    Chart/image-tag slides are left intact.
    """
    lines = [ln.rstrip() for ln in slide_md.split("\n") if ln.strip()]
    if not lines:
        return [slide_md]

    has_chart = any(parse_chart_tag(ln) for ln in lines)
    has_image = any(parse_image_tag(ln) for ln in lines)
    bullets = [ln for ln in lines if re.match(r"^\s*(?:[-*]|\d+\.)\s+", ln)]
    title_line = next((ln for ln in lines if ln.startswith("# ")), "# Slide")
    subtitle_line = next((ln for ln in lines if ln.startswith("## ")), "")
    type_line = next(
        (ln for ln in lines if re.match(r"^\[TYPE:", ln, re.IGNORECASE)),
        "[TYPE:text_heavy]",
    )

    avg_len = (sum(len(b) for b in bullets) / len(bullets)) if bullets else 0
    if has_chart or has_image or len(bullets) <= MAX_BULLETS_TWO_COL:
        return [slide_md]
    if len(bullets) <= MAX_BULLETS_TWO_COL + 2 and avg_len < 120:
        return [slide_md]

    split_at = max(MAX_BULLETS_SINGLE_COL, len(bullets) // 2)
    first_bullets = bullets[:split_at]
    second_bullets = bullets[split_at:]
    if not second_bullets:
        return [slide_md]

    first = [title_line]
    if subtitle_line:
        first.append(subtitle_line)
    first.append(type_line)
    first.extend(first_bullets)

    cont_title = title_line.replace("# ", "# ", 1)
    if not cont_title.lower().endswith("(cont.)"):
        cont_title = f"{cont_title} (cont.)"

    second = [cont_title, "## Continued"]
    second.append("[TYPE:text_heavy]")
    second.extend(second_bullets)

    return ["\n".join(first), "\n".join(second)]


def _looks_like_thank_you_slide(slide_md: str) -> bool:
    lines = [ln.strip() for ln in slide_md.splitlines() if ln.strip()]
    title_line = ""
    for ln in lines:
        if ln.startswith("# "):
            title_line = strip_md(ln)
            break
    text = f"{title_line} {' '.join(lines[:4])}".lower()
    return bool(re.search(r"thank\s*you|questions\??|q\s*&\s*a|closing", text))


# ═══════════════════════════════════════════════════════════════════════════════
# Chart PNG renderer via QuickChart.io  (port of renderChartPNG in TS)
# ═══════════════════════════════════════════════════════════════════════════════


async def render_chart_png(
    chart_type: str, labels: list[str], values: list[float]
) -> Optional[bytes]:
    """
    Render a chart to PNG via QuickChart.io (Chart.js v3 cloud renderer).
    Returns raw PNG bytes or None on failure.
    All chart_type strings match the TS equivalents:
      pie | doughnut | bar | col | line | area | radar | polar | stacked

    Visual theme: Slidesgo Sustainability & Environment — mint-green background,
    forest-green bars/lines, dark green labels, no heavy borders.
    """
    try:
        # Select palette deterministically from label content hash
        hash_val = sum(ord(c) for c in "".join(labels))
        pi = hash_val % len(PALETTES)
        pal = PALETTES[pi]
        primary = pal[0]
        pal_slice = pal[: len(values)]

        # Template-matched chart background (mint green slide bg)
        CHART_BG = "#EEF7E4"  # sage-green slide background—matches BG_LIGHT/CONTENT_BOX
        GRID_COLOR = "#C8E8A8"  # subtle sage-green grid lines
        TICK_COLOR = "#1E3A0E"  # dark forest green axis labels—matches HEADING_GREEN
        LEGEND_COLOR = "#1E3A0E"

        is_line_type = chart_type in ("line", "area")
        is_h_bar = chart_type == "bar"
        is_stacked = chart_type == "stacked"

        _fmt = "function(v){ if(Math.abs(v)>=1000000) return (v/1000000).toFixed(1)+'M'; if(Math.abs(v)>=1000) return (v/1000).toFixed(1)+'K'; return v; }"
        # _pct uses a reduce which QuickChart v3 doesn't reliably evaluate — replaced by _fmt_unit
        _fmt_unit = "function(v){ if(Math.abs(v)>=1000000) return (v/1000000).toFixed(1)+'M TCO\u2082Eq'; if(Math.abs(v)>=1000) return (v/1000).toFixed(1)+'K TCO\u2082Eq'; return parseFloat(v.toFixed(1))+' TCO\u2082Eq'; }"

        if chart_type == "doughnut":
            chart_config = {
                "type": "doughnut",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "data": values,
                            "backgroundColor": pal_slice,
                            "borderColor": CHART_BG,
                            "borderWidth": 3,
                            "hoverOffset": 10,
                        }
                    ],
                },
                "options": {
                    "cutout": "60%",
                    "plugins": {
                        "legend": {
                            "position": "bottom",
                            "labels": {
                                "color": LEGEND_COLOR,
                                "font": {"size": 28, "weight": "bold"},
                                "padding": 20,
                                "boxWidth": 24,
                            },
                        },
                        "datalabels": {
                            "color": "#ffffff",
                            "font": {"size": 22, "weight": "bold"},
                            "formatter": _fmt_unit,
                        },
                    },
                },
            }

        elif chart_type == "pie":
            chart_config = {
                "type": "pie",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "data": values,
                            "backgroundColor": pal_slice,
                            "borderColor": CHART_BG,
                            "borderWidth": 3,
                            "hoverOffset": 8,
                        }
                    ],
                },
                "options": {
                    "plugins": {
                        "legend": {
                            "position": "right",
                            "labels": {
                                "color": LEGEND_COLOR,
                                "font": {"size": 26, "weight": "bold"},
                                "padding": 16,
                                "boxWidth": 22,
                            },
                        },
                        "datalabels": {
                            "color": "#ffffff",
                            "font": {"size": 20, "weight": "bold"},
                            "formatter": _fmt_unit,
                        },
                    },
                },
            }

        elif chart_type == "polar":
            chart_config = {
                "type": "polarArea",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "data": values,
                            "backgroundColor": [_alpha_hex(c, 0.78) for c in pal_slice],
                            "borderColor": pal_slice,
                            "borderWidth": 2,
                        }
                    ],
                },
                "options": {
                    "plugins": {
                        "legend": {
                            "position": "bottom",
                            "labels": {
                                "color": LEGEND_COLOR,
                                "font": {"size": 26, "weight": "bold"},
                                "padding": 16,
                            },
                        },
                    },
                    "scales": {
                        "r": {
                            "grid": {"color": GRID_COLOR},
                            "ticks": {
                                "backdropColor": "transparent",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                                "stepSize": max(1, int(max(values) / 4))
                                if values
                                else 1,
                            },
                        },
                    },
                },
            }

        elif is_line_type:
            is_fill = chart_type == "area"
            chart_config = {
                "type": "line",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "label": "Value",
                            "data": values,
                            "borderColor": primary,
                            "backgroundColor": _alpha_hex(primary, 0.20)
                            if is_fill
                            else "transparent",
                            "fill": is_fill,
                            "tension": 0.40,
                            "pointRadius": 7,
                            "pointHoverRadius": 10,
                            "pointBackgroundColor": primary,
                            "pointBorderColor": CHART_BG,
                            "pointBorderWidth": 2,
                            "borderWidth": 3,
                        }
                    ],
                },
                "options": {
                    "plugins": {
                        "legend": {"display": False},
                        "datalabels": {
                            "color": TICK_COLOR,
                            "anchor": "end",
                            "align": "top",
                            "font": {"size": 26, "weight": "bold"},
                            "formatter": _fmt,
                        },
                    },
                    "layout": {
                        "padding": {"top": 38, "right": 12, "bottom": 4, "left": 4}
                    },
                    "scales": {
                        "y": {
                            "title": {
                                "display": True,
                                "text": "TCO\u2082Eq",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                            },
                            "grid": {"color": GRID_COLOR},
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxTicksLimit": 6,
                            },
                        },
                        "x": {
                            "grid": {"color": "transparent"},
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxRotation": 30,
                            },
                        },
                    },
                },
            }

        elif chart_type == "radar":
            chart_config = {
                "type": "radar",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "label": "Emissions",
                            "data": values,
                            "backgroundColor": _alpha_hex(primary, 0.28),
                            "borderColor": primary,
                            "borderWidth": 2.5,
                            "pointBackgroundColor": primary,
                            "pointBorderColor": CHART_BG,
                            "pointBorderWidth": 2,
                            "pointRadius": 5,
                        }
                    ],
                },
                "options": {
                    "plugins": {"legend": {"display": False}},
                    "scales": {
                        "r": {
                            "grid": {"color": GRID_COLOR},
                            "angleLines": {"color": GRID_COLOR},
                            "ticks": {
                                "backdropColor": "transparent",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                                "maxTicksLimit": 4,
                            },
                            "pointLabels": {
                                "color": TICK_COLOR,
                                "font": {"size": 28, "weight": "bold"},
                            },
                        },
                    },
                },
            }

        elif is_stacked:
            half = (len(labels) + 1) // 2
            pad_b = [0] * max(0, half - (len(values) - half))
            chart_config = {
                "type": "bar",
                "data": {
                    "labels": labels[:half],
                    "datasets": [
                        {
                            "label": "Primary",
                            "data": values[:half],
                            "backgroundColor": _alpha_hex(pal[0], 0.90),
                            "borderColor": pal[0],
                            "borderWidth": 0,
                            "borderRadius": 6,
                        },
                        {
                            "label": "Secondary",
                            "data": values[half : half * 2] + pad_b,
                            "backgroundColor": _alpha_hex(pal[1], 0.90),
                            "borderColor": pal[1],
                            "borderWidth": 0,
                            "borderRadius": 6,
                        },
                    ],
                },
                "options": {
                    "plugins": {
                        "legend": {
                            "position": "bottom",
                            "labels": {
                                "color": LEGEND_COLOR,
                                "font": {"size": 26, "weight": "bold"},
                                "boxWidth": 22,
                                "padding": 16,
                            },
                        },
                    },
                    "scales": {
                        "x": {
                            "stacked": True,
                            "grid": {"color": "transparent"},
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxRotation": 25,
                            },
                        },
                        "y": {
                            "stacked": True,
                            "title": {
                                "display": True,
                                "text": "TCO\u2082Eq",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                            },
                            "grid": {"color": GRID_COLOR},
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxTicksLimit": 5,
                            },
                        },
                    },
                },
            }

        else:
            # col (vertical bar) or bar (horizontal)
            bar_colors = [
                _alpha_hex(pal[i % len(pal)], 0.92) for i in range(len(values))
            ]
            bar_border = [pal[i % len(pal)] for i in range(len(values))]
            chart_config = {
                "type": "bar",
                "data": {
                    "labels": labels,
                    "datasets": [
                        {
                            "label": "tCO\u2082Eq",
                            "data": values,
                            "backgroundColor": bar_colors,
                            "borderColor": bar_border,
                            "borderWidth": 0,
                            "borderRadius": 6 if is_h_bar else 8,
                            "borderSkipped": False,
                        }
                    ],
                },
                "options": {
                    "indexAxis": "y" if is_h_bar else "x",
                    "plugins": {
                        "legend": {"display": False},
                        "datalabels": {
                            "color": TICK_COLOR,
                            "anchor": "end",
                            "align": "end" if is_h_bar else "top",
                            "font": {"size": 26, "weight": "bold"},
                            "formatter": _fmt,
                        },
                    },
                    "layout": {
                        "padding": {
                            "top": 4 if is_h_bar else 38,
                            "right": 110 if is_h_bar else 10,
                            "bottom": 4,
                            "left": 4,
                        }
                    },
                    "scales": {
                        "x": {
                            "grid": {
                                "color": GRID_COLOR if is_h_bar else "transparent"
                            },
                            "title": {
                                "display": is_h_bar,
                                "text": "TCO\u2082Eq",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                            },
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxRotation": 0 if is_h_bar else 30,
                            },
                        },
                        "y": {
                            "grid": {
                                "color": "transparent" if is_h_bar else GRID_COLOR
                            },
                            "title": {
                                "display": not is_h_bar,
                                "text": "TCO\u2082Eq",
                                "color": TICK_COLOR,
                                "font": {"size": 22},
                            },
                            "ticks": {
                                "color": TICK_COLOR,
                                "font": {"size": 26},
                                "maxTicksLimit": 6,
                            },
                        },
                    },
                },
            }

        payload = {
            "chart": chart_config,
            "width": 900,
            "height": 530,
            "backgroundColor": "transparent",
            "version": "3",
            "format": "png",
            "devicePixelRatio": 1,
        }
        async with httpx.AsyncClient(timeout=30) as client:
            res = await client.post("https://quickchart.io/chart", json=payload)
        if not res.is_success:
            print(f"QuickChart failed: {res.status_code}")
            return None
        return res.content

    except Exception as exc:
        print(f"render_chart_png error: {exc}")
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# Pexels image fetcher  (port of fetchPexelsImage in TS)
# ═══════════════════════════════════════════════════════════════════════════════

_ESG_KEYWORD_MAP: dict[str, str] = {
    "regulat": "esg regulation compliance policy green",
    "risk": "climate risk sustainability strategy",
    "opportunit": "green opportunity renewable energy future",
    "commit": "sustainability commitment green pledge net zero",
    "emiss": "carbon emissions clean energy environment",
    "scope": "greenhouse gas emissions carbon footprint",
    "waste": "waste management recycling sustainability",
    "water": "water conservation sustainability environment",
    "travel": "sustainable travel low carbon transport",
    "energy": "renewable energy solar wind clean power",
    "supply": "sustainable supply chain ESG procurement",
    "strateg": "sustainability strategy net zero roadmap",
    "target": "ESG targets net zero sustainability goals",
    "report": "ESG reporting sustainability transparency",
    "execut": "corporate sustainability leadership ESG",
    "kpi": "sustainability KPI metrics green performance",
    "intern": "global sustainability ESG international operations",
    "comparison": "sustainability comparison green benchmark",
}


def _esg_enrich_query(query: str) -> str:
    """Append ESG-relevant keywords to image search queries."""
    q_lower = query.lower()
    # If already has ESG/sustainability terms, just ensure 'sustainability' is present
    if any(
        w in q_lower
        for w in (
            "esg",
            "sustainability",
            "sustainable",
            "green",
            "carbon",
            "emission",
            "climate",
            "renewable",
        )
    ):
        return query + " sustainability"
    # Map known slide topic keywords to richer ESG queries
    for kw, enriched in _ESG_KEYWORD_MAP.items():
        if kw in q_lower:
            return enriched
    # Generic fallback: append sustainability context
    return query + " ESG sustainability environment green"


async def fetch_pexels_image(query: str) -> Optional[dict]:
    """
    Fetch a landscape stock photo from Pexels.
    Returns {"buffer": bytes, "mime": "image/jpeg"} or None.
    """
    query = _esg_enrich_query(query)

    async def _download_image(url: str) -> Optional[dict]:
        try:
            async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
                res = await client.get(url)
            if not res.is_success:
                return None
            content_type = (res.headers.get("content-type") or "image/jpeg").split(";")[
                0
            ]
            return {"buffer": res.content, "mime": content_type}
        except Exception:
            return None

    if not PEXELS_API_KEY or PEXELS_API_KEY == "your_pexels_api_key_here":
        print("PEXELS_API_KEY not configured — using public image fallback")
        q = quote_plus(query)
        fallback_urls = [
            f"https://source.unsplash.com/1600x900/?{q},sustainability",
            f"https://picsum.photos/seed/{q}/1600/900",
        ]
        for url in fallback_urls:
            img = await _download_image(url)
            if img:
                return img
        return None
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            search_res = await client.get(
                "https://api.pexels.com/v1/search",
                params={"query": query, "per_page": 3, "orientation": "landscape"},
                headers={"Authorization": PEXELS_API_KEY},
            )
        if not search_res.is_success:
            print(f"Pexels search failed: {search_res.status_code}, falling back")
            q = quote_plus(query)
            return await _download_image(f"https://picsum.photos/seed/{q}/1600/900")
        data = search_res.json()
        photo = (data.get("photos") or [None])[0]
        if not photo:
            q = quote_plus(query)
            return await _download_image(f"https://picsum.photos/seed/{q}/1600/900")
        img_url: str = (
            photo["src"].get("large2x")
            or photo["src"].get("large")
            or photo["src"].get("medium")
        )
        pexels_image = await _download_image(img_url)
        if pexels_image:
            return pexels_image
        q = quote_plus(query)
        return await _download_image(f"https://picsum.photos/seed/{q}/1600/900")
    except Exception as exc:
        print(f"fetch_pexels_image error: {exc}")
        q = quote_plus(query)
        return await _download_image(f"https://picsum.photos/seed/{q}/1600/900")


# ═══════════════════════════════════════════════════════════════════════════════
# Brand chrome — applied to every non-hero slide  (port of addSlideChrome in TS)
# ═══════════════════════════════════════════════════════════════════════════════

# ── Slide-type → template layout index mapping ───────────────────────────────
# The Slidesgo template typically carries these layouts:
#   0  Title Slide        (CENTER_TITLE + SUBTITLE)
#   1  Title + Content    (TITLE ph0 + BODY ph1 with styled bullet paragraphs)
#   2  Title + 2-Content  (TITLE + two side-by-side body columns)
#   3  Title + Content + Caption
#   4  Title Only
#   5  Two Content        (comparison/side-by-side)
#   6  Comparison         (labelled two-column)
#   7  Title + Picture
#   8  Blank
# We try multiple candidate indices in priority order and fall back gracefully.
_SLIDE_TYPE_LAYOUT_CANDIDATES: dict[str, list[int]] = {
    "hero": [0],
    "text_heavy": [1, 2, 3, 4],  # prefer body-placeholder layouts for template styling
    "summary": [3, 1, 4],
    "comparison": [5, 6, 2, 4],
    "chart_focus": [4, 1, 3],
    "image_focus": [10, 7, 1, 4],  # Prefer BLANK (10) to avoid placeholder conflicts
}


def _get_slide_layout(prs: Presentation, slide_type: str):
    """Return the most appropriate slide layout object for a given slide type.
    Prefers blank layouts to avoid placeholder conflicts and overlaps.
    """
    candidates = _SLIDE_TYPE_LAYOUT_CANDIDATES.get(slide_type, [10, 4, 1])
    for idx in candidates:
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            continue
    return prs.slide_layouts[min(10, len(prs.slide_layouts) - 1)]


def _clear_non_title_text_placeholders(slide) -> None:
    """
    Clear text from ALL non-title text-bearing placeholders so no default hint text
    (e.g. 'ANALYSIS A', 'Click to add text') leaks through from template layouts.
    Picture / media placeholders without text_frame are skipped.
    """
    for shp in list(slide.shapes):
        try:
            if not shp.is_placeholder:
                continue
            pf = shp.placeholder_format
            if pf.idx == 0:
                continue  # keep title placeholder intact
            # Only clear shapes that actually have a text frame
            if hasattr(shp, "text_frame"):
                shp.text_frame.clear()
        except Exception:
            continue


def _find_placeholder_by_type(slide, placeholder_type):
    """
    Find a placeholder by its type (TITLE, BODY, SUBTITLE, PICTURE, etc.)
    Returns the placeholder object or None if not found.
    """
    for shp in slide.shapes:
        try:
            if shp.is_placeholder and shp.placeholder_format.type == placeholder_type:
                return shp
        except Exception:
            continue
    return None


def _remove_subtitle_placeholder(slide) -> None:
    """Clear subtitle placeholder text so no default hint text appears."""
    for shp in list(slide.shapes):
        try:
            if not shp.is_placeholder:
                continue
            if shp.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                shp.text_frame.clear()
        except Exception:
            continue


def _clear_all_placeholders_except_title(slide) -> None:
    """
    Clear text from all placeholders except the title placeholder.
    This prevents hint text like 'Click to add text' from showing.
    """
    for shp in list(slide.shapes):
        try:
            if not shp.is_placeholder:
                continue
            # Skip title placeholder
            if shp.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                continue
            # Clear any placeholder with text frame
            if hasattr(shp, "text_frame") and shp.text_frame:
                shp.text_frame.clear()
        except Exception:
            continue


def _set_slide_title(slide, title: str, center_title: bool = False) -> None:
    """
    Fill the template layout's title placeholder with the slide title.
    Falls back to a textbox in the expected position if the placeholder is absent.
    """
    ph = _find_placeholder_by_type(slide, PP_PLACEHOLDER.TITLE)

    if ph:
        if not center_title:
            ph.left = Inches(0.35)
            ph.top = Inches(0.25)
            ph.width = Inches(9.3)
            ph.height = Inches(0.65)
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = title

        title_len = len(title)
        if title_len <= 40:
            title_font = 28
        elif title_len <= 60:
            title_font = 24
        elif title_len <= 80:
            title_font = 20
        else:
            title_font = 18

        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.LEFT if not center_title else PP_ALIGN.CENTER
            para.line_spacing = 1.0
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(title_font)
                run.font.color.rgb = hex_color(HEADING_GREEN)
                run.font.name = TEMPLATE_FONT_MAJOR

        _clear_all_placeholders_except_title(slide)
    else:
        # Fallback: draw a textbox where the template header area normally sits
        add_textbox(
            slide,
            title,
            0.35,
            0.25,
            9.3,
            0.6,
            font_size=20,
            bold=True,
            color_hex=HEADING_GREEN,
        )


def _fill_body_placeholder(slide, paragraphs: list[str], font_size: float = 14) -> bool:
    """
    Fill the slide's body placeholder with bullet paragraphs.
    This inherits ALL template styling — font family, color, spacing, bullets —
    defined in the slide layout and master.  Returns True if successful.
    """
    ph = _find_placeholder_by_type(slide, PP_PLACEHOLDER.BODY)
    if not ph:
        return False
    try:
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, text in enumerate(paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.level = 0
            run = para.add_run()
            run.text = text
            # Apply font size while keeping template's font face and color
            run.font.size = Pt(font_size)
        return True
    except Exception:
        return False


def _compute_body_top(title: str, subtitle: str) -> float:
    """Compute body start to avoid overlap with wrapped long titles."""
    title_len = len(title)
    body_y = 1.2 if not subtitle else 1.72
    if title_len > 55:
        body_y += 0.18
    if title_len > 75:
        body_y += 0.16
    return min(body_y, 2.2)


def _enrich_chart_bullets(
    labels: list[str], values: list[float], bullets: list[str]
) -> list[str]:
    """Guarantee chart slides have at least 3 readable insight/action bullets."""
    if len(bullets) >= 3:
        return bullets
    if not labels or not values or len(labels) != len(values):
        return bullets

    total = sum(values) if values else 0
    top_idx = max(range(len(values)), key=lambda i: values[i])
    low_idx = min(range(len(values)), key=lambda i: values[i])
    top_share = (values[top_idx] / total * 100) if total > 0 else 0
    low_share = (values[low_idx] / total * 100) if total > 0 else 0

    synthesized = [
        f"{labels[top_idx]} contributes {values[top_idx]:,.2f} tCO\u2082e ({top_share:.1f}% of total) \u2014 primary reduction focus.",
        f"Launch a 90-day abatement program for {labels[top_idx]} with monthly KPI tracking and owner accountability.",
        f"{labels[low_idx]} represents {values[low_idx]:,.2f} ({low_share:.1f}%) \u2014 maintain controls while prioritising high-impact sources.",
    ]

    merged = list(bullets)
    for b in synthesized:
        if len(merged) >= 4:
            break
        merged.append(b)
    return merged


# ═══════════════════════════════════════════════════════════════════════════════
# Per-slide renderer  (port of renderSlide() in TS)
# ═══════════════════════════════════════════════════════════════════════════════


async def render_slide(
    prs: Presentation,
    slide,
    slide_md: str,
    slide_num: int,
    total: int,
    logo_bytes: Optional[bytes],
    layout: Optional[dict],
):
    """
    Parse one slide's markdown and render all shapes onto the python-pptx slide.
    Supports layout types: hero, summary, comparison, chart_focus, image_focus, text_heavy.
    """
    # ── Clear ALL placeholders first to prevent template text leakage ──────────
    for shp in list(slide.shapes):
        try:
            if shp.is_placeholder and hasattr(shp, "text_frame") and shp.text_frame:
                shp.text_frame.clear()
        except Exception:
            pass

    # ── Parse slide markdown ──────────────────────────────────────────────────
    lines = [
        l.strip() for l in slide_md.split("\n") if l.strip() and "---SLIDE---" not in l
    ]

    title = f"Slide {slide_num}"
    subtitle = ""
    bullets: list[str] = []
    chart_tag: Optional[dict] = None
    image_tag: Optional[dict] = None

    for line in lines:
        if line.startswith("# "):
            title = strip_md(re.sub(r"^#+\s*", "", line))
        elif line.startswith("## "):
            subtitle = strip_md(re.sub(r"^#+\s*", "", line))
        elif (
            line.startswith("- ")
            or line.startswith("* ")
            or re.match(r"^\d+\.\s", line)
        ):
            raw = re.sub(r"^[-*]\s*", "", line)
            raw = re.sub(r"^\d+\.\s*", "", raw)
            # [A]/[B] comparison-column assignments — keep prefix intact for comparison renderer
            if re.match(r"^\[(?:A|B)\]", raw, re.IGNORECASE):
                bullets.append(raw)
            else:
                ct = parse_chart_tag(raw)
                it = parse_image_tag(raw)
                if ct and not chart_tag:
                    chart_tag = ct
                elif it and not image_tag:
                    image_tag = it
                elif not _is_dsl_tag(raw):
                    bullets.append(_clean_bullet_prefix(strip_md(raw)))
        elif not line.startswith("#"):
            # Always try to extract chart/image tags first, even from DSL-tagged lines
            ct = parse_chart_tag(line)
            it = parse_image_tag(line)
            if ct and not chart_tag:
                chart_tag = ct
            elif it and not image_tag:
                image_tag = it
            elif not _is_dsl_tag(line) and not subtitle:
                subtitle = strip_md(line)

    # Determine effective slide type
    layout_type = layout.get("slide_type") if layout else None
    # Normalise legacy / plan.json type names to the renderer's known types
    _TYPE_ALIASES: dict[str, str] = {
        "section": "hero",
        "kpi_template": "summary",
        "kpi": "summary",
        "three_column": "three_card",
        "two_column": "comparison",
        "four_column": "three_card",
        "news": "three_card",
        "bullets": "text_heavy",
        "chart": "chart_focus",
        "image": "image_focus",
        "title": "hero",
        "closing": "hero",
        "appendix": "text_heavy",
    }
    if layout_type:
        layout_type = _TYPE_ALIASES.get(layout_type, layout_type)

    # [TYPE:xxx] tag inside the slide markdown always wins — it is the AI's
    # explicit per-slide layout instruction and takes precedence over the plan.
    _inline_type_match = next(
        (
            re.match(r"^\[TYPE:([\w_]+)\]", ln, re.IGNORECASE)
            for ln in lines
            if re.match(r"^\[TYPE:", ln, re.IGNORECASE)
        ),
        None,
    )
    if _inline_type_match:
        _inline_type = _inline_type_match.group(1).lower()
        _inline_type = _TYPE_ALIASES.get(_inline_type, _inline_type)
        layout_type = _inline_type

    effective_type = layout_type or (
        "chart_focus" if chart_tag else "image_focus" if image_tag else "text_heavy"
    )

    # Last slide is always a clean closing/thank-you hero — override everything.
    if slide_num == total:
        effective_type = "hero"

    # ── HERO ──────────────────────────────────────────────────────────────────
    if effective_type == "hero":
        # Full-slide sage-green background.  Must be added BEFORE drawing text so
        # text textboxes (appended later = higher z-order) appear on top.
        add_rect(slide, 0.0, 0.0, W, H, BG_LIGHT)
        _add_slide_dot_clusters(slide)

        # Silence template CENTER_TITLE / SUBTITLE placeholders so their
        # dashed-border hint text doesn't show beneath our drawn content.
        for _ph in list(slide.placeholders):
            try:
                _ph.text_frame.clear()
            except Exception:
                pass

        # Hero title — drawn as a new textbox (on top of the BG_LIGHT rect)
        title_font = 36 if len(title) <= 40 else (30 if len(title) <= 65 else 24)
        add_textbox(
            slide,
            title,
            0.5,
            1.2,
            W - 1.0,
            2.0,
            font_size=title_font,
            bold=True,
            color_hex=HEADING_GREEN,
            align="center",
            valign="middle",
            wrap=True,
        )
        # Decorative green underline below the title
        add_rect(slide, W * 0.25, 3.35, W * 0.50, 0.04, BRAND_GREEN)

        # Subtitle / tagline
        hero_subtitle = subtitle
        if hero_subtitle and len(hero_subtitle) > 130:
            first_sent = re.split(r"[.!?]", hero_subtitle)[0].strip()
            hero_subtitle = (
                first_sent[:130] if len(first_sent) > 2 else hero_subtitle[:130]
            )
        if not hero_subtitle and bullets:
            hero_subtitle = bullets[0][:130]
        if hero_subtitle:
            add_textbox(
                slide,
                hero_subtitle,
                0.8,
                3.50,
                W - 1.6,
                1.0,
                font_size=15,
                color_hex=BODY_TEXT,
                italic=True,
                align="center",
                wrap=True,
            )
        return

    # ── SUMMARY / KPI CARDS ───────────────────────────────────────────────────
    if effective_type == "summary":
        _set_slide_title(slide, title, center_title=False)
        body_y = _compute_body_top(title, subtitle)
        body_h = H - body_y - 0.1
        add_rect(slide, 0.0, body_y - 0.05, W, H - body_y + 0.05, BG_LIGHT)
        _add_slide_dot_clusters(slide)

        # Detect STAT-format bullets: "STAT: Label | value | unit"
        stat_items: list[dict] = []
        non_stat_bullets: list[str] = []
        for b in bullets:
            st_match = re.match(
                r"^STAT:\s*([^|]+)\|\s*([^|]+)\|\s*(.+)$", b, re.IGNORECASE
            )
            pipe_match = (not st_match) and re.match(
                r"^([^|]+)\|\s*([^|]+)\|\s*(.+)$", b
            )
            if st_match:
                stat_items.append(
                    {
                        "label": st_match.group(1).strip(),
                        "value": st_match.group(2).strip(),
                        "unit": st_match.group(3).strip(),
                    }
                )
            elif pipe_match:
                stat_items.append(
                    {
                        "label": pipe_match.group(1).strip(),
                        "value": pipe_match.group(2).strip(),
                        "unit": pipe_match.group(3).strip(),
                    }
                )
            else:
                non_stat_bullets.append(b)

        # Slidesgo-style KPI cards: CONTENT_BOX fill, HEADING_GREEN label, ACCENT_ORANGE value
        # Matches the "Contents" table / card-grid style from the template screenshots.
        _kpi_icons = ["chart-bar", "factory", "building", "globe", "sun", "recycle"]

        if len(stat_items) >= 2:
            # KPI card grid (max 6 cards, up to 3 per row)
            per_row = len(stat_items) if len(stat_items) <= 3 else 3
            rows = (len(stat_items) + per_row - 1) // per_row
            card_w = (W - 0.3) / per_row - 0.12
            avail_card_h = H - 0.28 - body_y
            card_h = (
                min(2.2, avail_card_h - 0.1)
                if rows == 1
                else min(1.65, (avail_card_h - (rows - 1) * 0.12) / rows)
            )
            gap_x = card_w + 0.12
            gap_y = card_h + 0.12

            for idx, item in enumerate(stat_items[:6]):
                col = idx % per_row
                row = idx // per_row
                cx = 0.15 + col * gap_x
                cy = body_y + row * gap_y

                # Card: off-white sage fill + dark forest left border stripe
                add_rect(slide, cx, cy, card_w, card_h, CONTENT_BOX)
                add_rect(slide, cx, cy, 0.04, card_h, HEADING_GREEN)

                # Icon (top-left corner of card)
                icon_name = _kpi_icons[idx % len(_kpi_icons)]
                icon_bytes = render_icon_png(icon_name, 44, HEADING_GREEN)
                if icon_bytes:
                    add_picture_from_bytes(
                        slide, icon_bytes, cx + 0.10, cy + 0.08, 0.30, 0.30
                    )

                # Label (uppercase, dark forest green)
                add_textbox(
                    slide,
                    item["label"].upper(),
                    cx + 0.44,
                    cy + 0.08,
                    card_w - 0.54,
                    0.28,
                    font_size=7.5,
                    bold=True,
                    color_hex=HEADING_GREEN,
                )
                # Value (large, warm orange — primary focal element)
                v_font = 22 if len(item["value"]) > 8 else 28
                add_textbox(
                    slide,
                    item["value"],
                    cx + 0.08,
                    cy + 0.44,
                    card_w - 0.16,
                    card_h * 0.44,
                    font_size=v_font,
                    bold=True,
                    color_hex=ACCENT_ORANGE,
                    valign="middle",
                )
                # Unit (small caption, dark olive body text)
                add_textbox(
                    slide,
                    item["unit"],
                    cx + 0.08,
                    cy + card_h - 0.44,
                    card_w - 0.16,
                    0.40,
                    font_size=8.5,
                    color_hex=BODY_TEXT,
                    wrap=True,
                )

            # Non-stat bullets below cards
            if non_stat_bullets:
                below_y = body_y + rows * gap_y + 0.05
                add_bullet_textbox(
                    slide,
                    non_stat_bullets[:3],
                    0.15,
                    below_y,
                    W - 0.3,
                    H - 0.35 - below_y,
                    font_size=9.5,
                )
            # Add image at bottom of KPI summary for visual appeal only if there's space
            if rows == 1:
                summary_img = await fetch_pexels_image(f"{title} sustainability ESG")
                if summary_img:
                    img_w = 2.5
                    img_x = (W - img_w) / 2
                    img_y = H - 1.8
                    add_picture_from_bytes(
                        slide, summary_img["buffer"], img_x, img_y, img_w, 1.5
                    )
        else:
            # No KPI items - add an image with bullets
            summary_img = await fetch_pexels_image(f"{title} sustainability ESG")
            if summary_img:
                img_w = 4.0
                add_picture_from_bytes(
                    slide,
                    summary_img["buffer"],
                    W - img_w - 0.1,
                    body_y + 0.1,
                    img_w,
                    body_h - 0.2,
                )
                # Fallback: single column bullet layout with image
                add_bullet_textbox(
                    slide,
                    bullets[:5],
                    0.2,
                    body_y + 0.1,
                    W - img_w - 0.4,
                    body_h - 0.2,
                    font_size=11,
                )
            else:
                # Fallback: two-column bullet layout
                mid = (len(bullets) + 1) // 2
                col_w = (W - 0.55) / 2
                add_bullet_textbox(
                    slide,
                    bullets[:mid],
                    0.2,
                    body_y,
                    col_w,
                    H - body_y - 0.35,
                    font_size=10,
                )
                add_bullet_textbox(
                    slide,
                    bullets[mid:],
                    0.2 + col_w + 0.15,
                    body_y,
                    col_w,
                    H - body_y - 0.35,
                    font_size=10,
                )
        return

    # ── COMPARISON ────────────────────────────────────────────────────────────
    if effective_type == "comparison":
        _set_slide_title(slide, title, center_title=False)
        body_y = _compute_body_top(title, subtitle)
        body_h = H - body_y - 0.1
        add_rect(slide, 0.0, body_y - 0.05, W, body_h + 0.1, BG_LIGHT)
        _add_slide_dot_clusters(slide)

        col_w = 4.6
        col1_x, col2_x = 0.15, 5.1

        # Parse [A]/[B]-prefixed bullets or split evenly
        a_items: list[str] = []
        b_items: list[str] = []
        for b in bullets:
            if re.match(r"^\[A\]", b, re.IGNORECASE):
                a_items.append(re.sub(r"^\[A\]\s*", "", b, flags=re.IGNORECASE))
            elif re.match(r"^\[B\]", b, re.IGNORECASE):
                b_items.append(re.sub(r"^\[B\]\s*", "", b, flags=re.IGNORECASE))
            else:
                (a_items if len(a_items) <= len(b_items) else b_items).append(b)

        header_a = (
            "Current State"
            if a_items
            and re.search(
                r"current|existing|before|state|scope|emission|today",
                a_items[0],
                re.IGNORECASE,
            )
            else "Key Findings"
        )
        header_b = (
            "Recommendations"
            if b_items
            and re.search(
                r"target|future|after|recommend|goal|initiative|2030|next",
                b_items[0],
                re.IGNORECASE,
            )
            else "Target State"
        )

        # Determine icon per card based on header/title content
        def _pick_comparison_icon(header: str, pos: int) -> str:
            h = (header + " " + title).lower()
            if any(
                w in h
                for w in [
                    "factory",
                    "combustion",
                    "scope 1",
                    "stationary",
                    "fuel",
                    "diesel",
                ]
            ):
                return "factory"
            if any(
                w in h
                for w in [
                    "scope 3",
                    "travel",
                    "supply",
                    "global",
                    "international",
                    "commute",
                ]
            ):
                return "globe"
            if any(
                w in h
                for w in [
                    "electric",
                    "bolt",
                    "scope 2",
                    "power",
                    "energy",
                    "renewable",
                    "solar",
                ]
            ):
                return "bolt"
            if any(
                w in h
                for w in [
                    "recommend",
                    "target",
                    "goal",
                    "initiative",
                    "2030",
                    "roadmap",
                ]
            ):
                return "award"
            if any(
                w in h for w in ["current", "existing", "finding", "analysis", "state"]
            ):
                return "building"
            if any(w in h for w in ["waste", "recycle", "circular"]):
                return "recycle"
            if any(w in h for w in ["water", "liquid", "supply"]):
                return "water"
            if any(w in h for w in ["people", "employee", "staff", "users", "commute"]):
                return "users"
            return ["factory", "globe", "trending-up", "award", "leaf", "chart-bar"][
                pos % 6
            ]

        icon_a = _pick_comparison_icon(header_a, 0)
        icon_b = _pick_comparison_icon(header_b, 1)

        # Icon-above-text card layout (matching "In depth analysis" template slide)
        for card_i, (col_x, items, header, icon_name) in enumerate(
            [
                (col1_x, a_items, header_a, icon_a),
                (col2_x, b_items, header_b, icon_b),
            ]
        ):
            # Card background (off-white sage green)
            add_rect(slide, col_x, body_y, col_w, body_h, CONTENT_BOX)
            # Bottom accent line
            add_rect(slide, col_x, body_y + body_h - 0.04, col_w, 0.04, BRAND_GREEN)

            # Icon centred at top of card (orange, matching template)
            icon_bytes = render_icon_png(icon_name, 72, ACCENT_ORANGE)
            if icon_bytes:
                icon_x = col_x + (col_w - 0.58) / 2
                add_picture_from_bytes(
                    slide, icon_bytes, icon_x, body_y + 0.14, 0.58, 0.58
                )

            # Card heading (bold dark forest green, centred)
            add_textbox(
                slide,
                header,
                col_x + 0.10,
                body_y + 0.80,
                col_w - 0.20,
                0.40,
                font_size=13,
                bold=True,
                color_hex=HEADING_GREEN,
                align="center",
            )
            # Green underline beneath heading
            add_rect(
                slide,
                col_x + col_w * 0.25,
                body_y + 1.22,
                col_w * 0.50,
                0.025,
                BRAND_GREEN,
            )

            # Bullet items (dark olive, centred)
            content_items = (
                items if items else (bullets[:4] if card_i == 0 else bullets[4:8])
            )
            if content_items:
                # Dynamically size font based on content length to avoid overflow
                total_chars = sum(len(c) for c in content_items[:5])
                dynamic_font = 9.5 if total_chars > 350 else (10.5 if total_chars > 250 else 11.5)
                
                add_bullet_textbox(
                    slide,
                    content_items[:5],
                    col_x + 0.14,
                    body_y + 1.28,
                    col_w - 0.28,
                    body_h - 1.36,
                    font_size=dynamic_font,
                    color_hex=BODY_TEXT,
                )

        # Add image to comparison slide for visual appeal
        comp_img = await fetch_pexels_image(f"{title} sustainability comparison")
        if comp_img and not (a_items and b_items):
            # If no comparison content, show image full width
            img_h = min(body_h * 0.5, 2.0)
            add_picture_from_bytes(
                slide, comp_img["buffer"], 0.15, body_y + 0.15, W - 0.3, img_h
            )

        # Fallback: if both columns empty, render two-column text_heavy
        if not a_items and not b_items and bullets:
            mid = (len(bullets) + 1) // 2
            add_bullet_textbox(
                slide,
                bullets[:mid],
                col1_x + 0.14,
                body_y + 1.28,
                col_w - 0.28,
                body_h - 1.36,
                font_size=11.5,
            )
            add_bullet_textbox(
                slide,
                bullets[mid:],
                col2_x + 0.14,
                body_y + 1.28,
                col_w - 0.28,
                body_h - 1.36,
                font_size=11.5,
            )
        return

    # ── THREE_CARD (3-column icon-card grid, inspired by "News examples" template slide) ─
    if effective_type == "three_card":
        _set_slide_title(slide, title, center_title=False)
        body_y = _compute_body_top(title, subtitle)
        body_h = H - body_y - 0.1
        add_rect(slide, 0.0, body_y - 0.05, W, body_h + 0.1, BG_LIGHT)
        _add_slide_dot_clusters(slide)

        # Parse [A]/[B]/[C]-prefixed bullets into groups
        groups: list[list[str]] = [[], [], [], []]
        for b in bullets:
            m = re.match(r"^\[([ABCD])\]\s*(.*)", b, re.IGNORECASE)
            if m:
                idx = ord(m.group(1).upper()) - ord("A")
                if idx < 4:
                    groups[idx].append(m.group(2).strip())
            else:
                # Distribute evenly if no prefix
                smallest = min(range(3), key=lambda i: len(groups[i]))
                groups[smallest].append(b)

        # Determine number of cards (2 or 3) and assign icons
        active_groups = [g for g in groups[:3] if g]
        n_cards = max(len(active_groups), 2 if bullets else 0)
        if n_cards == 0:
            n_cards = 3

        # Fallback: if no prefix groups, redistribute bullets into n_cards equal slices
        if not any(groups[:3]) and bullets:
            chunk = max(1, len(bullets) // 3)
            groups[0] = bullets[:chunk]
            groups[1] = bullets[chunk : chunk * 2]
            groups[2] = bullets[chunk * 2 :]

        _card_icons = ["bolt", "recycle", "shield", "leaf", "award", "trending-up"]
        _card_titles = ["New Project", "Key Initiative", "CSR & Impact"]

        card_gap = 0.12
        card_w = (W - 0.30 - card_gap * 2) / 3
        card_h = body_h

        for ci in range(3):
            cx = 0.15 + ci * (card_w + card_gap)
            cy = body_y

            # Card background
            add_rect(slide, cx, cy, card_w, card_h, CONTENT_BOX)
            # Thin bottom accent line
            add_rect(slide, cx, cy + card_h - 0.04, card_w, 0.04, BRAND_GREEN)

            # Icon centred at top (orange)
            ic_name = _card_icons[ci % len(_card_icons)]
            ic_bytes = render_icon_png(ic_name, 68, ACCENT_ORANGE)
            if ic_bytes:
                ix = cx + (card_w - 0.54) / 2
                add_picture_from_bytes(slide, ic_bytes, ix, cy + 0.14, 0.54, 0.54)

            # Derive card title from first bullet or index
            card_title = _card_titles[ci] if not groups[ci] else ""
            if groups[ci]:
                # Use first item as heading if it looks like a title (short)
                if len(groups[ci][0]) < 55:
                    card_title = groups[ci][0]
                    card_bullets = groups[ci][1:]
                else:
                    card_title = _card_titles[ci]
                    card_bullets = groups[ci]
            else:
                card_bullets = []

            add_textbox(
                slide,
                card_title,
                cx + 0.08,
                cy + 0.78,
                card_w - 0.16,
                0.40,
                font_size=11,
                bold=True,
                color_hex=HEADING_GREEN,
                align="center",
            )
            add_rect(
                slide, cx + card_w * 0.25, cy + 1.20, card_w * 0.50, 0.022, BRAND_GREEN
            )

            if card_bullets:
                add_bullet_textbox(
                    slide,
                    card_bullets[:4],
                    cx + 0.10,
                    cy + 1.26,
                    card_w - 0.20,
                    card_h - 1.34,
                    font_size=10.5,
                    color_hex=BODY_TEXT,
                )
        return

    # ── Shared chrome for remaining types (chart_focus, image_focus, text_heavy) ─────────
    # Draw body background FIRST so that title and all content textboxes appended
    # afterwards sit on top (higher z-order = visible).
    body_y = _compute_body_top(title, subtitle)
    body_h = H - body_y - 0.1
    add_rect(slide, 0.0, body_y - 0.05, W, body_h + 0.1, BG_LIGHT)
    _add_slide_dot_clusters(slide)
    # Title is drawn via add_textbox so it ends up ABOVE the BG_LIGHT rect.
    # _set_slide_title still fills the placeholder (useful if template rendering
    # is active), but the visible label comes from the textbox appended here.
    _set_slide_title(slide, title, center_title=False)
    # Subtitle subheading — only when bullets also exist
    subtitle_as_subheading = bool(subtitle and bullets)
    if subtitle_as_subheading:
        add_textbox(
            slide,
            subtitle,
            0.35,
            body_y,
            W - 0.7,
            0.44,
            font_size=12,
            italic=True,
            color_hex=HEADING_GREEN,
            align="center",
            wrap=True,
        )
        body_y += 0.52
        body_h -= 0.52

    # ── CHART_FOCUS ───────────────────────────────────────────────────────────
    if effective_type == "chart_focus" and chart_tag:
        chart_zone = (
            (layout or {}).get("layout_metadata", {}).get("chart_zone", "left 55%")
        )
        pct_match = re.search(r"(\d+)%", chart_zone)
        chart_width_pct = min(
            float(pct_match.group(1)) / 100 if pct_match else 0.50, 0.52
        )
        c_w = (W - 0.25) * chart_width_pct
        c_h = min(body_h, 3.1)
        chart_off_y = body_y + (body_h - c_h) / 2  # vertically centred
        text_x = c_w + 0.28
        text_w = W - text_x - 0.1

        chart_buffer = await render_chart_png(
            chart_tag["type"], chart_tag["labels"], chart_tag["values"]
        )
        if chart_buffer:
            # Primary path: QuickChart PNG embedded as picture
            add_picture_from_bytes(slide, chart_buffer, 0.12, chart_off_y, c_w, c_h)
        else:
            # Fallback path: native python-pptx chart (bar or pie only)
            _type_map = {
                "col": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "area": XL_CHART_TYPE.AREA,
                "pie": XL_CHART_TYPE.PIE,
                "doughnut": XL_CHART_TYPE.DOUGHNUT,
                "stacked": XL_CHART_TYPE.COLUMN_STACKED,
            }
            xl_type = _type_map.get(chart_tag["type"], XL_CHART_TYPE.COLUMN_CLUSTERED)
            chart_data = ChartData()
            chart_data.categories = chart_tag["labels"]
            chart_data.add_series("Value", chart_tag["values"])
            slide.shapes.add_chart(
                xl_type,
                Inches(0.12),
                Inches(chart_off_y),
                Inches(c_w),
                Inches(c_h),
                chart_data,
            )

        # Only show non-STAT bullets in the text panel (STAT data is already in the chart)
        visible_bullets = [
            b
            for b in bullets
            if not re.match(r"^STAT:\s*", b, re.IGNORECASE)
            and not re.match(r"^[^|]+\|[^|]+\|.+$", b)
        ]
        visible_bullets = _enrich_chart_bullets(
            chart_tag["labels"], chart_tag["values"], visible_bullets
        )
        if text_w > 1.5:
            # Coloured legend swatches matching QuickChart's palette (Slidesgo template style)
            pal_idx = sum(ord(c) for c in "".join(chart_tag["labels"])) % len(PALETTES)
            pal = PALETTES[pal_idx]
            legend_n = min(len(chart_tag["labels"]), 7)
            legend_item_h = min(0.45, body_h / max(legend_n + 1, 1))
            # Unit suffix for all emissions chart types (including pie/doughnut/polar)
            _unit_sfx = " TCO\u2082Eq"
            for li in range(legend_n):
                ly = body_y + li * legend_item_h
                swatch_color = pal[li % len(pal)].lstrip("#")
                add_rect(
                    slide,
                    text_x,
                    ly + legend_item_h * 0.18,
                    0.14,
                    legend_item_h * 0.64,
                    swatch_color,
                )
                lbl = chart_tag["labels"][li]
                v_val = chart_tag["values"][li]
                v_str = (
                    f"{v_val:,.1f}{_unit_sfx}"
                    if isinstance(v_val, (int, float))
                    else str(v_val)
                )
                add_textbox(
                    slide,
                    lbl,
                    text_x + 0.20,
                    ly,
                    text_w - 0.22,
                    legend_item_h * 0.52,
                    font_size=11,
                    bold=True,
                    color_hex=HEADING_GREEN,
                    valign="bottom",
                )
                add_textbox(
                    slide,
                    v_str,
                    text_x + 0.20,
                    ly + legend_item_h * 0.50,
                    text_w - 0.22,
                    legend_item_h * 0.50,
                    font_size=10,
                    color_hex=BODY_TEXT,
                    valign="top",
                )
            # Insight bullets below legend if space remains
            if visible_bullets:
                below_y = body_y + legend_n * legend_item_h + 0.12
                rem_h = body_y + body_h - below_y
                if rem_h > 0.45:
                    add_bullet_textbox(
                        slide,
                        visible_bullets[:2],
                        text_x,
                        below_y,
                        text_w,
                        rem_h,
                        font_size=10,
                        color_hex=BODY_TEXT,
                    )
        return
    # chart_focus without a chart tag — fall through to text_heavy

    # ── IMAGE_FOCUS ───────────────────────────────────────────────────────────
    if effective_type == "image_focus" or image_tag:
        image_query = (image_tag or {}).get("query") or f"{title} sustainability"
        image_layout = (image_tag or {}).get("layout") or "right"
        # Reformat any STAT-format bullets to readable text
        img_bullets = (
            _reformat_stat_bullets(bullets[:MAX_BULLETS_TWO_COL]) if bullets else []
        )
        image_payload = await fetch_pexels_image(image_query)

        if image_payload and image_layout == "bg":
            add_picture_from_bytes(
                slide, image_payload["buffer"], 0.0, body_y, W, body_h
            )
            add_rect(slide, 0.16, body_y + 0.18, W * 0.58, body_h - 0.36, "FFFFFF")
            if img_bullets:
                add_bullet_textbox(
                    slide,
                    img_bullets,
                    0.24,
                    body_y + 0.28,
                    W * 0.54,
                    body_h - 0.56,
                    font_size=13,
                )
        elif image_payload and image_layout == "left":
            # Use smaller image to avoid overlap with text
            img_w = 3.8
            img_h = body_h * 0.7  # Only 70% of body height to leave room for text
            add_picture_from_bytes(
                slide, image_payload["buffer"], 0.12, body_y + body_h * 0.15, img_w, img_h
            )
            # Only show bullets, not title - title should only be in header
            if img_bullets:
                text_x = img_w + 0.25
                text_w = W - text_x - 0.15
                add_bullet_textbox(
                    slide, img_bullets, text_x, body_y, text_w, body_h, font_size=12
                )
            # Don't show title as content - it's only in the header
        elif image_payload:
            # Use smaller image to avoid overlap
            img_w = 3.8
            img_h = body_h * 0.7  # Only 70% of body height
            add_picture_from_bytes(
                slide, image_payload["buffer"], W - img_w - 0.12, body_y + body_h * 0.15, img_w, img_h
            )
            # Only show bullets, not title - title should only be in header
            if img_bullets:
                add_bullet_textbox(
                    slide,
                    img_bullets,
                    0.18,
                    body_y,
                    W - img_w - 0.4,
                    body_h,
                    font_size=12,
                )
            # Don't show title as content - it's only in the header
            else:
                add_textbox(
                    slide,
                    text_content[0],
                    0.22,
                    body_y + 0.3,
                    W - img_w - 0.54,
                    body_h - 0.6,
                    font_size=17,
                    color_hex=BODY_TEXT,
                    align="left",
                    valign="middle",
                    wrap=True,
                )
        elif img_bullets:
            # Show bullets but no image if image fetch failed
            add_bullet_textbox(
                slide, img_bullets, 0.18, body_y, W - 0.3, body_h, font_size=13
            )
        # If no bullets and no image, leave content empty - title is in header only
        return

    # ── TEXT_HEAVY / default ──────────────────────────────────────────────────
    if bullets:
        # Convert any STAT-format bullets ('STAT: label | value | unit') to readable text
        all_bullets = _reformat_stat_bullets(bullets[:MAX_BULLETS_TWO_COL])

        # Always use manual textboxes - never rely on placeholders to avoid overlap issues
        # For slides with ≤4 bullets, try to add a right-side image
        # Determine if we should add an image based on slide content
        # Aim for 4-5 images in a 15-slide deck
        add_image = len(all_bullets) <= 5 and len(all_bullets) >= 2
        img_payload = None
        if add_image:
            img_payload = await fetch_pexels_image(f"{title} sustainability")

        # Calculate optimal font size based on bullet count and available space
        bullet_count = len(all_bullets)
        if bullet_count <= 3:
            base_font_size = 14
        elif bullet_count <= 5:
            base_font_size = 12
        elif bullet_count <= 7:
            base_font_size = 10.5
        else:
            base_font_size = 9.5

        if img_payload and len(all_bullets) <= 4:
            # Right-image layout: bullets left, photo right (smaller image to avoid overlap)
            img_w = 3.5
            img_h = body_h * 0.75  # Smaller height to leave room for bullets
            add_picture_from_bytes(
                slide, img_payload["buffer"], W - img_w - 0.12, body_y, img_w, img_h
            )
            add_bullet_textbox(
                slide,
                all_bullets,
                0.18,
                body_y,
                W - img_w - 0.42,
                body_h,
                font_size=base_font_size,
            )
        elif len(all_bullets) >= 5:
            # Three-column layout for many bullets to prevent overflow
            cols = 3
            bullets_per_col = (bullet_count + cols - 1) // cols
            col_w = (W - 0.4) / cols
            for col_idx in range(cols):
                start_idx = col_idx * bullets_per_col
                end_idx = min(start_idx + bullets_per_col, bullet_count)
                col_bullets = all_bullets[start_idx:end_idx]
                if col_bullets:
                    col_x = 0.15 + col_idx * (col_w + 0.05)
                    add_bullet_textbox(
                        slide,
                        col_bullets,
                        col_x,
                        body_y,
                        col_w - 0.05,
                        body_h,
                        font_size=base_font_size - 0.5,
                    )
        elif len(all_bullets) >= 3:
            # Two-column layout
            mid = (len(all_bullets) + 1) // 2
            col_w = (W - 0.55) / 2
            add_bullet_textbox(
                slide,
                all_bullets[:mid],
                0.2,
                body_y,
                col_w,
                body_h,
                font_size=base_font_size,
            )
            add_bullet_textbox(
                slide,
                all_bullets[mid:],
                0.2 + col_w + 0.15,
                body_y,
                col_w,
                body_h,
                font_size=base_font_size,
            )
            divx = 0.2 + col_w + 0.075
            add_rect(slide, divx, body_y + 0.1, 0.008, body_h - 0.2, "DDDDDD")
        elif len(all_bullets) <= 2:
            # Use manual textbox for small number of bullets
            # Slidesgo-style statement cards: CONTENT_BOX bg + icon + accent bar + BODY_TEXT
            _stmt_icons = [
                "trending-up",
                "award",
                "leaf",
                "shield",
                "globe",
                "chart-bar",
            ]
            card_h = min(
                (body_h - 0.1 * (len(all_bullets) - 1)) / max(len(all_bullets), 1),
                1.6,
            )
            for i, b in enumerate(all_bullets):
                cy = body_y + i * (card_h + 0.14)
                # Card background
                add_rect(slide, 0.18, cy, W - 0.36, card_h, CONTENT_BOX)
                # Left accent bar
                add_rect(slide, 0.18, cy, 0.05, card_h, BRAND_GREEN)
                # Icon
                ic_bytes = render_icon_png(
                    _stmt_icons[i % len(_stmt_icons)], 44, ACCENT_ORANGE
                )
                if ic_bytes:
                    add_picture_from_bytes(
                        slide, ic_bytes, 0.28, cy + (card_h - 0.34) / 2, 0.34, 0.34
                    )
                # Text
                add_textbox(
                    slide,
                    b,
                    0.72,
                    cy,
                    W - 0.90,
                    card_h,
                    font_size=15,
                    color_hex=BODY_TEXT,
                    align="left",
                    valign="middle",
                    wrap=True,
                )
        else:
            # Always use manual textbox - never rely on placeholders
            add_bullet_textbox(
                slide, all_bullets, 0.2, body_y, W - 0.35, body_h, font_size=14
            )
    else:
        # No bullets at all — just show a small image if available, no text content
        # Title should only be in header, not in body content
        img_pl = await fetch_pexels_image(f"{title} sustainability")
        if img_pl:
            img_w = 3.5
            img_h = body_h * 0.7  # Smaller to leave empty space
            add_picture_from_bytes(
                slide, img_pl["buffer"], W - img_w - 0.12, body_y, img_w, img_h
            )
        # If no image either, leave content empty - title is only in header


# ═══════════════════════════════════════════════════════════════════════════════
# Chutes LLM caller  (port of callChutes in TS)
# ═══════════════════════════════════════════════════════════════════════════════


async def call_chutes(prompt: str, max_tokens: int = 4000) -> str:
    """
    Call Chutes API, cascading through models until one succeeds.
    Mirrors the model-fallback loop in route.ts.
    """
    models = [
        "Qwen/Qwen3-235B-A22B-Instruct-2507-TEE",
        "deepseek-ai/DeepSeek-R1-Distill-Llama-70B",
        "openai/gpt-oss-20b",
        "chutesai/Mistral-Small-3.2-24B-Instruct-2506",
    ]
    last_status = 0
    if not CHUTES_API_KEY:
        raise RuntimeError(
            "CHUTES_API_KEY environment variable is not set. "
            "Add it before starting the server."
        )
    async with httpx.AsyncClient(timeout=120) as client:
        for model in models:
            try:
                res = await client.post(
                    CHUTES_URL,
                    headers={
                        "Content-Type": "application/json",
                        "Authorization": f"Bearer {CHUTES_API_KEY}",
                    },
                    json={
                        "model": model,
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.7,
                        "max_tokens": max_tokens,
                    },
                )
                if res.is_success:
                    print(f"Chutes: using model {model}")
                    data = res.json()
                    return data["choices"][0]["message"]["content"] or ""
                last_status = res.status_code
                print(f"Model {model} failed with {res.status_code}, trying next...")
            except Exception as exc:
                print(f"Model {model} exception: {exc}, trying next...")
    raise RuntimeError(f"All Chutes models failed. Last status: {last_status}")


# ═══════════════════════════════════════════════════════════════════════════════
# AI layout planner  (port of generateLayoutPlan in TS)
# ═══════════════════════════════════════════════════════════════════════════════


async def generate_layout_plan(
    topic: str, context_summary: str, slide_count: int
) -> list[dict]:
    """
    Ask the LLM to produce a JSON layout plan for the entire deck.
    Falls back to a fixed-pattern plan on any error.
    """
    prompt = f"""You are an expert slide designer. Given a topic and dataset, produce a JSON layout plan for the entire deck.

Topic: {topic}
Data summary: {context_summary[:800]}
Slide count: {slide_count}
Style: modern, professional, bold visual hierarchy, balanced whitespace

JSON schema for each object in the array:
{{
  "slide_index": number,
  "slide_type": "hero" | "text_heavy" | "chart_focus" | "comparison" | "image_focus" | "summary",
  "layout_metadata": {{
    "text_zone": "e.g. left 45% | top_center 60%",
    "chart_zone": "e.g. left 65% | center 70% (omit if no chart)",
    "image_zone": "e.g. right 50% | bg 100% (omit if no image)",
    "design_style": "minimal | bold | professional | creative"
  }},
  "content_intent": "One sentence: what this slide communicates",
  "chart_intent": "Chart type and data fields if chart is expected (omit if no chart)",
  "image_keywords": ["keyword1", "keyword2"]
}}

Rules:
- Slide 1 is always hero (title slide)
- Last slide is hero (closing/thank-you)
- Use summary type for KPI/metrics overview slides (3-6 big numbers)
- Use comparison type for side-by-side analysis slides
- Use chart_focus for data-heavy slides — at least 35% of slides should have charts
- Use image_focus for narrative/context slides to add visual impact
- Vary slide_type across the deck — no two consecutive same type unless necessary

IMPORTANT: Output ONLY valid JSON array. No markdown, no explanations, no text before or after. Start with [ and end with ]."""

    try:
        raw = await call_chutes(prompt, 2000)
        print(f"Layout plan raw response (first 500 chars): {raw[:500]}")

        # Try multiple extraction strategies
        json_str = None
        errors = []

        # Strategy 1: Extract from markdown code block
        for pattern in [
            r"```json\s*(\[[\s\S]*?\])\s*```",
            r"```\s*(\[[\s\S]*?\])\s*```",
        ]:
            match = re.search(pattern, raw, re.IGNORECASE)
            if match:
                json_str = match.group(1)
                break

        # Strategy 2: Find JSON array with balanced brackets
        if not json_str:
            # Find the first [ and try to get a balanced JSON array
            start_idx = raw.find("[")
            if start_idx != -1:
                # Try to find matching ] by counting brackets
                depth = 0
                in_string = False
                escape_next = False
                for i, char in enumerate(raw[start_idx:], start_idx):
                    if escape_next:
                        escape_next = False
                        continue
                    if char == "\\":
                        escape_next = True
                        continue
                    if char == '"' and not escape_next:
                        in_string = not in_string
                        continue
                    if in_string:
                        continue
                    if char == "[":
                        depth += 1
                    elif char == "]":
                        depth -= 1
                        if depth == 0:
                            json_str = raw[start_idx : i + 1]
                            break

        # Strategy 3: Simple regex as last resort
        if not json_str:
            match = re.search(r"\[[\s\S]*\]", raw)
            if match:
                json_str = match.group(0)

        if not json_str:
            raise ValueError("No JSON array found in layout plan response")

        # Clean up the JSON string
        json_str = json_str.strip()

        # Try to parse, if fails, try to fix common issues
        try:
            plan: list[dict] = json.loads(json_str)
        except json.JSONDecodeError as e:
            # Try to fix common JSON issues
            print(f"Initial JSON parse failed: {e}, attempting fixes...")
            # Remove trailing commas
            json_str = re.sub(r",\s*]", "]", json_str)
            json_str = re.sub(r",\s*}", "}", json_str)
            # Remove markdown
            json_str = re.sub(r"^```json\s*", "", json_str, flags=re.IGNORECASE)
            json_str = re.sub(r"```$", "", json_str, flags=re.IGNORECASE)
            plan = json.loads(json_str)

        print(f"Layout plan: {len(plan)} slides planned")
        return plan
    except Exception as err:
        print(f"Layout plan generation failed, using default plan: {err}")
        types = [
            "hero",
            "summary",
            "chart_focus",
            "chart_focus",
            "image_focus",
            "text_heavy",
            "comparison",
            "chart_focus",
            "image_focus",
            "chart_focus",
            "text_heavy",
            "chart_focus",
            "summary",
            "comparison",
            "hero",
        ]
        return [
            {
                "slide_index": i,
                "slide_type": types[i % len(types)],
                "layout_metadata": {"design_style": "professional"},
            }
            for i in range(slide_count)
        ]


# ═══════════════════════════════════════════════════════════════════════════════
# Main AI slide generator  (port of generateSlides in TS)
# ═══════════════════════════════════════════════════════════════════════════════

_FORMAT_RULES = """STRICT OUTPUT FORMAT — follow exactly:
- Put this exact delimiter alone on its own line between every slide: ---SLIDE---
- Every slide MUST start with: # Slide Title
- Optionally one ## subtitle sentence after the title
- Then 4-6 bullet points starting with "- "
- NO markdown bold (**), NO extra headers, plain professional English
- Do NOT write [SLIDE N] or slide-number markers anywhere — they corrupt the output

SLIDE TYPE TAGS — include exactly ONE per slide on its own line after ##:
[TYPE:hero]          — section-break / intro slide, no bullet list needed, just a title + subtitle
[TYPE:summary]       — KPI/metrics overview; bullets MUST use format: - STAT: Label | value | unit
                       Example: - STAT: Scope 3 Total | 45,840 TCO2Eq | 52.3% of portfolio
[TYPE:chart_focus]   — data-heavy slide; include a [CHART:...] tag; chart takes 65% of slide
[TYPE:comparison]    — side-by-side analysis; prefix bullets with [A] or [B] to assign to each column
[TYPE:image_focus]   — narrative/visual slide; include an [IMAGE:...] tag for a Pexels stock photo
[TYPE:text_heavy]    — dense analysis slide, auto-splits into two columns

CHART TAG SYNTAX:  [CHART:type:Label1=value1,Label2=value2,...]
  Types: pie | doughnut | bar (horizontal) | col (vertical bar) | line | area | radar | polar | stacked
IMAGE TAG SYNTAX (only on slides WITHOUT a chart): [IMAGE:search query:layout]
  Layouts: bg | right | left | inset

All numeric bullets must include: real number from data + what it means + specific action/implication.
Do NOT prefix bullet text with labels like 'ACTION:', 'INSIGHT:', 'RECOMMENDATION:', 'KEY POINT:', 'FINDING:' — write plain, professional bullet sentences directly.
No text before first ---SLIDE--- and no text after last slide."""


def _parse_slide_chunks(raw: str, label: str) -> list[str]:
    """Split raw LLM output into individual slide markdown strings."""
    chunks = [
        s.strip() for s in re.split(r"---+\s*slide\s*---+", raw, flags=re.IGNORECASE)
    ]
    chunks = [c for c in chunks if c and "#" in c]
    if len(chunks) < 2:
        print(
            f"[{label}] delimiter split yielded {len(chunks)} — falling back to heading split"
        )
        by_heading = [s.strip() for s in re.split(r"(?=^# )", raw, flags=re.MULTILINE)]
        by_heading = [s for s in by_heading if s.startswith("#") and len(s) > 3]
        if len(by_heading) > len(chunks):
            chunks = by_heading
    print(f"[{label}] parsed {len(chunks)} slides")
    return chunks


async def generate_slides(context_data: str, page_type: str) -> dict:
    """
    Run AI layout planning and content generation in parallel.
    Returns {"slides_markdown": list[str], "layout_plan": list[dict]}.
    Mirrors generateSlides() in route.ts.
    """
    is_report = page_type == "ghg-report"
    label = "GHG Emissions Report" if is_report else "ESG Dashboard"
    slide_count = 15 if is_report else 8

    # Begin layout plan generation (runs concurrently with slide content)
    layout_plan_task = asyncio.create_task(
        generate_layout_plan(label, context_data, slide_count)
    )

    if not is_report:
        # ── 8-slide dashboard — single AI call ────────────────────────────────
        prompt = f"""You are a senior ESG consultant creating an {label} for Growlity's Board of Directors.
{_FORMAT_RULES}

Create 8 compelling slides using varied types: at least 1 hero, 1 summary, 3 chart_focus, 1 comparison, 1 image_focus.
Cover: title (hero), KPI overview (summary), scope emissions breakdown, key category analysis, risks, strategy, recommendations.

DATA:
{context_data}

Generate all 8 slides now:"""
        raw, layout_plan = await asyncio.gather(
            call_chutes(prompt, 4500),
            layout_plan_task,
        )
        print(f"[Dashboard] raw output (first 600 chars): {raw[:600]}")
        slides_markdown = _parse_slide_chunks(raw, "Dashboard")
        print(f"Dashboard: AI generated {len(slides_markdown)} slides")
        return {"slides_markdown": slides_markdown, "layout_plan": layout_plan}

    # ── 15-slide GHG report — two parallel AI calls ───────────────────────────
    shared_ctx = f"DATA (use exact numbers for all bullets and charts):\n{context_data}"

    prompt_part1 = f"""You are a senior ESG consultant creating a GHG Emissions Report for Growlity's Board of Directors.
{_FORMAT_RULES}

Generate EXACTLY slides 1 through 8 of a 15-slide deck.
Slide type variety required: slide 1 hero, slide 2 summary (KPI cards), slides 3-7 mix of chart_focus and text_heavy, slide 8 comparison or image_focus.
Cover: title/overview (hero), KPI metrics (summary), Scope 1 three sub-categories (chart_focus), Scope 2 energy (chart_focus), Scope 3 overview (chart_focus), plus 1 comparison and 1 image slide.

{shared_ctx}

Output ONLY slides 1-8 (each separated by ---SLIDE---):"""

    prompt_part2 = f"""You are a senior ESG consultant creating a GHG Emissions Report for Growlity's Board of Directors.
{_FORMAT_RULES}

Generate EXACTLY slides 9 through 15 of a 15-slide deck.
Slide type variety required: at least 2 chart_focus, 1 comparison, 1 image_focus, 1 summary or hero intro, 1 text_heavy, last slide hero (thank-you).
Cover: Scope 3 deep-dives (business travel chart, logistics, water/waste/people), top emitters ranked (bar chart), regulatory risks (comparison or text_heavy), final recommendations (summary), closing (hero).

{shared_ctx}

Output ONLY slides 9-15 (each separated by ---SLIDE---):"""

    print(
        "GHG Report: making 3 parallel AI calls (layout plan + slides 1-8 + slides 9-15)..."
    )
    raw1, raw2, layout_plan = await asyncio.gather(
        call_chutes(prompt_part1, 4500),
        call_chutes(prompt_part2, 4500),
        layout_plan_task,
    )

    part1 = _parse_slide_chunks(raw1, "part1")
    part2 = _parse_slide_chunks(raw2, "part2")
    slides_markdown = part1 + part2
    print(
        f"GHG Report: part1={len(part1)}, part2={len(part2)}, total={len(slides_markdown)}"
    )
    return {"slides_markdown": slides_markdown, "layout_plan": layout_plan}


# ═══════════════════════════════════════════════════════════════════════════════
# FastAPI application  (equivalent to the Next.js POST handler in route.ts)
# ═══════════════════════════════════════════════════════════════════════════════

app = FastAPI(title="ESGtech.ai PPT Generator", version="1.0.0")


class PptRequest(BaseModel):
    contextData: str
    pageType: str = "dashboard"


@app.post("/generate-ppt")
async def generate_ppt(body: PptRequest):
    """
    POST /generate-ppt
    Body: {"contextData": "<text>", "pageType": "dashboard" | "ghg-report"}
    Returns: application/vnd.openxmlformats-officedocument.presentationml.presentation
    """
    if not body.contextData.strip():
        from fastapi import HTTPException

        raise HTTPException(status_code=400, detail="Missing contextData in request.")

    # ── Load logo bytes ────────────────────────────────────────────────────────
    logo_bytes: Optional[bytes] = None
    try:
        logo_bytes = LOGO_PATH.read_bytes()
    except FileNotFoundError:
        print(f"Logo not found at: {LOGO_PATH}")

    # ── AI generates slide content + layout plan ───────────────────────────────
    try:
        result = await generate_slides(body.contextData, body.pageType)
        slides_markdown: list[str] = result["slides_markdown"]
        layout_plan: list[dict] = result["layout_plan"]
        if not slides_markdown:
            raise ValueError("No slides returned")

        expanded_slides: list[str] = []
        expanded_plan: list[dict] = []
        for i, slide_md in enumerate(slides_markdown):
            lp = layout_plan[i] if i < len(layout_plan) else {}
            parts = _split_dense_slide_markdown(slide_md)
            for part_idx, part in enumerate(parts):
                expanded_slides.append(part)
                plan_entry = dict(lp) if lp else {}
                if part_idx > 0:
                    plan_entry["slide_type"] = "text_heavy"
                expanded_plan.append(plan_entry)

        slides_markdown = expanded_slides
        layout_plan = expanded_plan
    except Exception as err:
        print(f"AI slide generation failed: {err}")
        from fastapi import HTTPException

        raise HTTPException(
            status_code=503,
            detail="AI slide generation failed. Please try again in a moment.",
        )

    # ── Build PPTX — load Slidesgo template as the base ──────────────────────
    if _TEMPLATE_PATH.exists():
        prs = Presentation(str(_TEMPLATE_PATH))
        print(f"Using template: {_TEMPLATE_PATH.name}")
        # Strip sample slides — keep masters/layouts/theme/fonts only
        _clear_template_slides(prs)
    else:
        print(f"Template not found at {_TEMPLATE_PATH} — using blank presentation")
        prs = Presentation()

    # Ensure 16:9 dimensions (10 × 5.625 inches)
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # Use the TITLE layout (idx 0) for hero/title slides so CENTER_TITLE and SUBTITLE
    # placeholders are available.  Content slides use a type-appropriate layout so the
    # template's infographic decorations, body placeholder styles, and icon shapes are
    # preserved on each slide.
    try:
        layout_title = prs.slide_layouts[_TEMPLATE_TITLE_LAYOUT_IDX]
    except IndexError:
        layout_title = prs.slide_layouts[0]

    total = len(slides_markdown)

    for i, slide_md in enumerate(slides_markdown):
        # Determine layout before adding the slide so the template master is correct
        lp_entry = layout_plan[i] if i < len(layout_plan) else {}
        declared_type = (lp_entry or {}).get("slide_type", "")
        # Only treat as hero if: first slide, OR last slide with no real bullet content
        # and no explicit content type declared in the layout plan.
        _lines_with_content = [
            l.strip()
            for l in slide_md.splitlines()
            if l.strip() and not l.strip().startswith("#")
        ]
        _is_content_slide = (
            declared_type
            in ("comparison", "chart_focus", "image_focus", "text_heavy", "summary")
            or len(_lines_with_content) > 3
        )
        is_hero = (
            i == 0 or i == total - 1  # last slide is always a closing/thank-you hero
        )
        # Resolve the effective slide type early (using the same alias map as render_slide)
        # so we can pick the right template layout before adding the slide.
        _raw_type = (lp_entry or {}).get("slide_type", "")
        _TYPE_ALIASES_LOCAL: dict[str, str] = {
            "section": "hero",
            "kpi_template": "summary",
            "kpi": "summary",
            "three_column": "three_card",
            "two_column": "comparison",
            "four_column": "three_card",
            "news": "three_card",
            "bullets": "text_heavy",
            "chart": "chart_focus",
            "image": "image_focus",
            "title": "hero",
            "closing": "hero",
            "appendix": "text_heavy",
        }
        _eff_type = _TYPE_ALIASES_LOCAL.get(_raw_type, _raw_type) or "text_heavy"
        chosen_layout = layout_title if is_hero else _get_slide_layout(prs, _eff_type)
        slide = prs.slides.add_slide(chosen_layout)

        await render_slide(
            prs, slide, slide_md, i + 1, total, logo_bytes, lp_entry or None
        )

    # ── Serialise PPTX to bytes and return ────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    out_name = f"growlity-esg-report-{timestamp}.pptx"
    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{out_name}"',
            "X-PPT-Renderer": "python-ghg-v2",
            "X-PPT-Filename": out_name,
        },
    )


# ── Entry point for direct execution ─────────────────────────────────────────
if __name__ == "__main__":
    import uvicorn

    uvicorn.run("generate_ppt:app", host="0.0.0.0", port=8000, reload=True)
